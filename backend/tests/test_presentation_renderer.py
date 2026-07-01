"""Unit-тесты PresentationRenderer (layout router, overflow, KPI, gov-badge)."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from app.agents.models import AnalysisResult, AskResult
from app.presentation_renderer import (
    BADGE_TEXT,
    LayoutKind,
    PresentationRenderer,
    PresentationTheme,
)
from core.models import ChartSpec


def _ask(
    *,
    chart_type: str = "bar",
    x: str = "region",
    y: str = "debt",
    data: list[dict] | None = None,
    insights: list[str] | None = None,
) -> AskResult:
    return AskResult(
        question="Тестовый вопрос",
        sql="SELECT 1",
        data=data
        or [{"region": "г. Минск", "debt": 100}, {"region": "Гомельская область", "debt": 200}],
        analysis=AnalysisResult(
            insights=insights or ["инсайт 1", "инсайт 2", "инсайт 3"],
            key_conclusion="Главный вывод",
            reasoning="mock",
        ),
        chart_spec=ChartSpec(
            chart_type=chart_type,  # type: ignore[arg-type]
            title="Тест",
            x=x,
            y=y,
            rationale="mock",
        ),
        reasoning="mock",
    )


def test_choose_layout_router():
    r = PresentationRenderer()
    assert r.choose_layout("line") == LayoutKind.TOP_BOTTOM
    assert r.choose_layout("waterfall") == LayoutKind.TOP_BOTTOM
    assert r.choose_layout("donut") == LayoutKind.SIDE_BY_SIDE
    assert r.choose_layout("treemap") == LayoutKind.SIDE_BY_SIDE
    assert r.choose_layout("horizontal_bar") == LayoutKind.CHART_TABLE
    assert r.choose_layout("bar") == LayoutKind.CHART_TABLE
    assert r.choose_layout("kpi") == LayoutKind.KPI_FULL
    # scatter группируется с bar-семейством -> CHART_TABLE (актуальный роутер)
    assert r.choose_layout("scatter") == LayoutKind.CHART_TABLE


def test_truncate_text_overflow():
    r = PresentationRenderer()
    long = "а" * 200
    out = r.truncate_text(long, max_chars=120)
    assert len(out) == 120
    assert out.endswith("...")


def test_resolve_headers_action_title():
    r = PresentationRenderer()
    res = _ask()
    res.chart_spec = ChartSpec(
        chart_type="horizontal_bar",
        title="Заголовок",
        x="region",
        y="debt",
        action_title="Гомельская область — лидер",
        rationale="r",
    )
    header, subtitle = r._resolve_headers("Исходный вопрос?", res)
    assert header == "Гомельская область — лидер"
    assert subtitle == "Исходный вопрос?"


def test_build_ranking_table_rows_top5():
    r = PresentationRenderer()
    res = _ask(
        data=[
            {"region": "A", "debt": 10},
            {"region": "B", "debt": 50},
            {"region": "C", "debt": 30},
            {"region": "D", "debt": 20},
            {"region": "E", "debt": 40},
            {"region": "F", "debt": 5},
        ]
    )
    rows = r._build_ranking_table_rows(res, res.chart_spec)
    assert len(rows) == 5
    assert rows[0][0] == "B"
    # Денежные колонки форматируются с суффиксом "бел. руб."
    assert "бел. руб." in rows[0][1]


def test_gov_badge_on_title_slide(tmp_path):
    prs = Presentation()
    prs.slide_width = Inches(PresentationTheme.SLIDE_W)
    prs.slide_height = Inches(PresentationTheme.SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    PresentationRenderer().create_title_slide(slide)
    texts = "".join(sh.text for sh in slide.shapes if hasattr(sh, "text"))
    assert BADGE_TEXT in texts
    assert "BI-аналитика налогов РБ" in texts
    out = tmp_path / "test_title.pptx"
    prs.save(str(out))
    assert out.exists()


def test_chart_table_layout_adds_table(tmp_path):
    prs = Presentation()
    prs.slide_width = Inches(PresentationTheme.SLIDE_W)
    prs.slide_height = Inches(PresentationTheme.SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    renderer = PresentationRenderer()
    res = _ask(chart_type="horizontal_bar")
    renderer.render_slide_header(slide, "Топ регионов", "Вопрос?")
    # Рейтинговая таблица (правая часть CHART_TABLE layout) добавляет pptx-таблицу
    renderer._render_ranking_table(
        slide,
        res,
        res.chart_spec,
        left=0.5,
        top=2.0,
        width=4.0,
    )
    has_table = any(getattr(sh, "has_table", False) for sh in slide.shapes)
    assert has_table


def test_summary_slide_kpi_cards(tmp_path):
    prs = Presentation()
    prs.slide_width = Inches(PresentationTheme.SLIDE_W)
    prs.slide_height = Inches(PresentationTheme.SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    res = _ask(
        data=[
            {"region": "г. Минск", "tax_type": "НДС", "accrued": 100, "debt": 50},
            {"region": "Гомельская область", "tax_type": "Подоходный", "accrued": 80, "debt": 40},
        ]
    )
    PresentationRenderer().create_summary_slide(slide, "Обзор презентации.", res)
    texts = "".join(sh.text for sh in slide.shapes if hasattr(sh, "text"))
    # Summary-слайд содержит блок ключевых параметров и обзор
    assert "КЛЮЧЕВЫЕ ПАРАМЕТРЫ" in texts or "АНАЛИТИЧЕСКИЙ ОБЗОР" in texts
    assert BADGE_TEXT in texts
