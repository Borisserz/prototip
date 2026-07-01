import logging
import time
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pydantic import BaseModel, Field
from pypdf import PdfReader

from app.presentation_renderer import PresentationRenderer, PresentationTheme
from core import storage
from core.llm import call_structured

logger = logging.getLogger("PdfPresentation")


class SlideContent(BaseModel):
    title: str = Field(..., description="Заголовок слайда (короткий)")
    bullets: list[str] = Field(..., description="3-5 ключевых мыслей для слайда")


class DocumentPresentationSpec(BaseModel):
    overview: str = Field(..., description="Краткое описание документа (2-3 предложения)")
    themes: list[str] = Field(..., description="Главные темы документа")
    slides: list[SlideContent] = Field(..., description="Слайды с контентом")
    takeaways: list[str] = Field(..., description="Главные выводы из документа")


def generate_presentation_from_pdf(pdf_path: str) -> str:
    """Парсит PDF, прогоняет через LLM, возвращает путь к .pptx."""
    start_time = time.time()

    # 1. Извлечение текста
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"

    if not text.strip():
        raise ValueError("Не удалось извлечь текст из PDF")

    # Ограничиваем текст для LLM (чтобы влезло в контекст)
    text_snippet = text[:15000]

    logger.info(f"PDF прочитан: извлечено {len(text)} символов.")

    # 2. Вызов LLM для структуры
    prompt = f"""Ты — аналитик. Изучи текст документа и создай структуру презентации на русском языке.
Текст документа:
{text_snippet}
"""
    spec = call_structured(
        prompt,
        schema=DocumentPresentationSpec,
        system="Сгенерируй структуру презентации из текста. Отвечай только JSON.",
    )

    logger.info(f"LLM сгенерировала структуру: {len(spec.slides)} слайдов.")

    # 3. Генерация .pptx (напрямую через renderer)
    out_dir = Path("out")
    out_dir.mkdir(parents=True, exist_ok=True)
    presentation_id = uuid.uuid4().hex[:8]
    pptx_path = out_dir / f"doc_presentation_{presentation_id}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(PresentationTheme.SLIDE_W)
    prs.slide_height = Inches(PresentationTheme.SLIDE_H)
    blank = prs.slide_layouts[6]

    renderer = PresentationRenderer()

    # Титульный
    renderer.create_title_slide(prs.slides.add_slide(blank))

    # Обзор и темы
    renderer.create_summary_slide(prs.slides.add_slide(blank), spec.overview, None)

    # Слайды с контентом (используем takeaways slide layout для текста)
    for slide_data in spec.slides:
        slide = prs.slides.add_slide(blank)
        renderer._add_title_text(slide, slide_data.title, top=0.30, font_size=28)

        tx = slide.shapes.add_textbox(
            Inches(renderer.t.MARGIN_H), Inches(1.0), Inches(renderer.t.USABLE_W), Inches(5.0)
        )
        tf = tx.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(slide_data.bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {renderer.truncate_text(bullet, 150)}"
            p.font.size = renderer.t.Pt(18)
            p.font.name = renderer.t.FONT
            p.space_before = renderer.t.Pt(10)

        renderer._add_gov_badge(slide)
        renderer._add_footer(slide)

    # Выводы
    renderer.create_takeaways_slide(prs.slides.add_slide(blank), spec.takeaways)

    prs.save(str(pptx_path))
    # зеркалим .pptx в MinIO (неломающе)
    storage.mirror_artifact(pptx_path, "presentations")
    elapsed = int((time.time() - start_time) * 1000)
    logger.info(f"PDF Presentation готова: {pptx_path} ({elapsed}ms)")

    return str(pptx_path)
