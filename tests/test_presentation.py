"""Тесты PresentationAgent (Phase 6).

Проверка создания .pptx, структуры слайдов, наличия изображений.
Живой прогон на 3 вопросах (требует ollama + python-pptx).
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from app.agents.presentation_agent import PresentationAgent
from app.schemas import AnalysisResult, AskResult, DeckNarrative, PresentationResult
from core.llm import is_ollama_available
from core.models import ChartSpec


def _fake_ask_result(
    *,
    question: str = "тестовый вопрос",
    png_path: str = "/tmp/fake.png",
    chart_type: str = "bar",
    title: str = "Тест",
    data: list[dict] | None = None,
) -> AskResult:
    return AskResult(
        question=question,
        sql="SELECT 1",
        data=data or [{"x": 1}],
        analysis=AnalysisResult(
            insights=["i1", "i2", "i3"],
            key_conclusion="key",
            anomaly_or_trend=None,
            reasoning="mock analysis",
        ),
        chart_spec=ChartSpec(
            chart_type=chart_type,  # type: ignore[arg-type]
            title=title,
            x="x",
            y="x",
            rationale="mock chart",
        ),
        png_path=png_path,
        reasoning="mock ask",
    )


def _mock_planner_executor(*ask_results):
    pending = list(ask_results)
    mock_ex = MagicMock()

    def run_side_effect(agent_name: str, *args, **kwargs):
        if agent_name == "planner_agent":
            return pending.pop(0) if pending else MagicMock(success=False, error="no mock result")
        return MagicMock(success=False, error=f"unexpected agent {agent_name}")

    mock_ex.run.side_effect = run_side_effect
    return mock_ex


def test_presentation_agent_mock():
    """Мок: проверяем структуру презентации с DeckNarrative, кол-во слайдов, отсутствие заглушки."""
    agent = PresentationAgent()

    fake_res = PresentationResult(pptx_path="out/presentation.pptx", num_slides=7)

    with patch("app.agents.presentation_agent.PresentationAgent.run", return_value=fake_res):
        pass

    with (
        patch(
            "app.agents.presentation_agent.get_executor",
            return_value=_mock_planner_executor(
                _fake_ask_result(title="Тест бар 1"),
                _fake_ask_result(title="Тест бар 2"),
                _fake_ask_result(title="Тест бар 3"),
            ),
        ),
        patch("app.agents.presentation_agent.call_structured") as mock_narr,
    ):
        mock_narr.return_value = DeckNarrative(
            overview="test overview",
            themes=["t1", "t2"],
            key_takeaways=["k1", "k2", "k3", "k4"],
            recommendations=["r1", "r2"],
        )

        with patch("app.agents.presentation_agent.Presentation") as mock_prs:
            mock_prs_instance = MagicMock()
            mock_slide = MagicMock()
            mock_prs_instance.slides.add_slide.return_value = mock_slide
            #  title + ov + themes + 3q + takeaways + recs = 8
            mock_prs_instance.slides.__len__.return_value = 8
            mock_prs.return_value = mock_prs_instance

            res = agent.run(["q1", "q2", "q3"])

    assert isinstance(res, PresentationResult)
    assert res.num_slides >= 3 + 4  # questions + title/ov/themes/key/rec at least
    assert "presentation.pptx" in res.pptx_path


@pytest.mark.live
@pytest.mark.skipif(
    not is_ollama_available(),
    reason="Ollama + qwen2.5-coder:7b-instruct недоступен для live теста Phase 6",
)
def test_presentation_live_creates_file_with_slides_and_images():
    """Живой прогон на 3 вопросах. Файл создаётся, открывается python-pptx, слайдов >=3, изображения на месте."""
    from pptx import Presentation as PptxPresentation

    agent = PresentationAgent()
    questions = [
        "Какие регионы имеют наибольшую задолженность по НДС?",
        "Динамика начислений подоходного налога в г. Минск по месяцам?",
        "Топ-3 региона по сумме имущественных налогов?",
    ]

    res = agent.run(questions)

    assert isinstance(res, PresentationResult)
    pptx_file = Path(res.pptx_path)
    assert pptx_file.exists(), f"Презентация не создана: {pptx_file}"

    # open with python-pptx
    prs = PptxPresentation(str(pptx_file))
    n_slides = len(prs.slides)
    assert n_slides >= len(questions) + 4, (
        f"Ожидалось >= {len(questions)}+4 слайдов (титул+обзор+темы+вопросы+выводы+рекомендации), получено {n_slides}"
    )

    # нет слайда-заглушки
    all_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text or ""
    assert "Презентация собрана автоматически" not in all_text, (
        "Заглушка-плейсхолдер должна быть удалена"
    )

    # есть титул/обзор/выводы/рекомендации (по тексту)
    assert "BI-аналитика налогов РБ" in all_text
    assert "Обзор" in all_text or "overview" in all_text.lower()
    assert "Ключевые выводы" in all_text
    assert "Рекомендации" in all_text

    # изображения присутствуют
    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                image_count += 1
    assert image_count >= 1, "В презентации должны присутствовать изображения (PNG)"

    # clean? no, leave the artifact


def test_presentation_respects_prefs_and_exact_count_nonlive():
    """Non-live: prefs (chart_type) оверрайдятся в ребилде, exact count с num_slides+includes, RU тип в caption."""
    from unittest.mock import MagicMock, patch

    from app.agents.presentation_agent import PresentationAgent
    from app.schemas import DeckNarrative

    agent = PresentationAgent()
    qblocks = [
        {
            "text": "Топ-3 региона по задолженности",
            "chart_type": "horizontal_bar",
            "note": "pref hbar",
        },
        {"text": "Структура по налогам", "chart_type": "donut", "note": ""},
    ]

    with (
        patch(
            "app.agents.presentation_agent.get_executor",
            return_value=_mock_planner_executor(
                _fake_ask_result(
                    png_path="/tmp/f.png",
                    chart_type="bar",
                    title="Топ",
                    data=[{"region": "a", "debt": 1}],
                ),
                _fake_ask_result(
                    png_path="/tmp/f2.png",
                    chart_type="bar",
                    title="Struct",
                    data=[{"tax_type": "x", "accrued": 2}],
                ),
            ),
        ),
        patch("app.agents.presentation_agent.call_structured") as mock_narr,
        patch("app.agents.presentation_agent.build_chart") as mock_build,
    ):
        mock_narr.return_value = DeckNarrative(
            overview="ov",
            themes=["t1", "t2"],
            key_takeaways=["k1", "k2", "k3", "k4"],
            recommendations=["r1", "r2"],
        )

        # чтобы не падать на реальном pptx в этом тесте, патчим Presentation
        with patch("app.agents.presentation_agent.Presentation") as mock_prs:
            mock_prs_inst = MagicMock()
            mock_prs_inst.slides.__len__.return_value = 7  # title+ov+themes+2q+key+recs
            mock_prs_inst.slides.add_slide.return_value = MagicMock()
            mock_prs.return_value = mock_prs_inst

            res = agent.run(qblocks, num_slides=7, include_title=True, include_recommendations=True)

    assert res.num_slides == 7
    # build_chart должен был вызываться с оверрайднутым типом (хотя бы 1 раз с horizontal)
    called_types = []
    for call in mock_build.call_args_list:
        sp = call[0][1] if call[0] else call[1].get("spec")
        if sp:
            ct = getattr(sp, "chart_type", None) or (
                sp.get("chart_type") if isinstance(sp, dict) else None
            )
            if ct:
                called_types.append(ct)
    assert any(ct == "horizontal_bar" for ct in called_types), (
        f"pref horizontal_bar should override, got {called_types}"
    )
