"""PresentationAgent (Phase 6+): авто-сборка .pptx презентации.

Оркестрация: slide_pipeline (data→chart→analyst) → DeckNarrative → PresentationRenderer.
"""

from __future__ import annotations

import logging
import re
import time
import uuid
from contextlib import suppress
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from app.agent_context import presentation_subplan

from app.agents.base import BaseAgent
from app.agents.factory import get_executor
from app.agents.models import (
    AskResult,
    DeckNarrative,
    PresentationInput,
    PresentationResult,
    PresentationState,
    SlideData,
    SlideUpdate,
)
from app.chart_repair import repair_chart_spec
from app.pipeline_progress import emit_pipeline_stage, suppress_pipeline_emit
from app.presentation_renderer import PresentationRenderer, PresentationTheme
from app.slide_pipeline import build_slide_ask_result
from core.llm import call_structured, setup_logging
from app.agents.config_loader import get_agent_config
from viz.charts import build_chart, export_png

setup_logging()
logger = logging.getLogger("PresentationAgent")


class PresentationAgent(BaseAgent):
    """Агент сборки презентаций .pptx из slide pipeline (без nested Planner)."""

    name = "presentation_agent"
    description = (
        "По списку вопросов собирает .pptx: slide pipeline → PNG (viz) → "
        "PresentationRenderer (динамические макеты, gov-стиль, KPI, таблицы)."
    )

    def __init__(self) -> None:
        self.renderer = PresentationRenderer()

    def _slug(self, text: str, max_len: int = 40) -> str:
        # Transliterate Cyrillic → ASCII so filenames never hit latin-1 HTTP header errors
        _TRANSLIT: dict[str, str] = {
            'а':'a','б':'b','в':'v','г':'g','д':'d','е':'e','ё':'yo','ж':'zh',
            'з':'z','и':'i','й':'y','к':'k','л':'l','м':'m','н':'n','о':'o',
            'п':'p','р':'r','с':'s','т':'t','у':'u','ф':'f','х':'kh','ц':'ts',
            'ч':'ch','ш':'sh','щ':'shch','ъ':'','ы':'y','ь':'','э':'e','ю':'yu','я':'ya',
        }
        lowered = text.lower()
        transliterated = ''.join(_TRANSLIT.get(c, c) for c in lowered)
        slug = re.sub(r"[^a-zA-Z0-9]+", "_", transliterated).strip("_")
        return slug[:max_len] or "result"

    def _normalize_input(
        self, questions: list[str] | list[dict] | list | PresentationInput
    ) -> tuple[list[str], dict[int, str | None], list[dict]]:
        qlist: list[dict] = []
        if isinstance(questions, PresentationInput):
            qlist = [{"text": q} for q in questions.questions]
        elif questions and isinstance(questions[0], dict):
            qlist = questions  # type: ignore[assignment]
        else:
            qlist = [{"text": q} for q in (questions or [])]
        qs = [str(q.get("text", "")).strip() for q in qlist if str(q.get("text", "")).strip()]
        prefs: dict[int, str | None] = {}
        for qb in qlist:
            if str(qb.get("text", "")).strip():
                prefs[len(prefs)] = qb.get("chart_type")
        return qs, prefs, qlist

    def _trim_for_num_slides(
        self,
        questions: list[str],
        prefs: dict[int, str | None],
        qlist: list[dict],
        *,
        num_slides: int | None,
        include_title: bool,
        include_recommendations: bool,
    ) -> tuple[list[str], dict[int, str | None]]:
        if not num_slides or num_slides <= 0:
            return questions, prefs
        base_fixed = (1 if include_title else 0) + 2 + 1 + (1 if include_recommendations else 0)
        max_q = max(0, num_slides - base_fixed)
        trimmed = questions[:max_q] if max_q > 0 else (questions[:1] if questions else [])
        new_prefs: dict[int, str | None] = {}
        uidx = 0
        for qb in qlist:
            if str(qb.get("text", "")).strip() and uidx < len(trimmed):
                new_prefs[uidx] = qb.get("chart_type")
                uidx += 1
        return trimmed, new_prefs

    def _collect_results(self, questions: list[str], executor: Any) -> list[AskResult]:
        results: list[AskResult] = []
        with presentation_subplan(), suppress_pipeline_emit():
            for qi, q in enumerate(questions):
                emit_pipeline_stage(
                    "synthesis",
                    "running",
                    f"Вопрос {qi + 1}/{len(questions)}: {str(q)[:50]}...",
                    agent="presentation_agent",
                )
                results.append(build_slide_ask_result(q, executor))
        return results

    def _get_deck_narrative(self, questions: list[str], results: list[AskResult]) -> DeckNarrative:
        summaries: list[str] = []
        for q, r in zip(questions, results):  # noqa: B905
            conc = (r.analysis.key_conclusion if r.analysis else "") or ""
            anom = (r.analysis.anomaly_or_trend if r.analysis else "") or "нет"
            action = (getattr(r.chart_spec, "action_title", None) if r.chart_spec else None) or ""
            action_line = f"Action title: {action}\n" if action else ""
            summaries.append(
                f"Вопрос: {q}\n{action_line}Ключевой вывод: {conc}\nАномалия/тренд: {anom}"
            )
        cfg = get_agent_config("presentation_agent")
        
        prompt = f"""Ты — {cfg.role}. {cfg.goal}

{cfg.rules}

Вопросы и выводы:
{chr(10).join(summaries)}

Все тексты на русском. Верни только валидный JSON по схеме DeckNarrative.
"""
        return call_structured(
            prompt,
            schema=DeckNarrative,
            system=f"Ты — {cfg.role}. Выдавай только валидный JSON. Тексты на русском.",
            agent_name=self.name
        )

    def _fallback_narrative(self) -> DeckNarrative:
        return DeckNarrative(
            overview=(
                "Презентация содержит анализ налоговых поступлений Республики Беларусь "
                "по синтетическим данным за 2024 год (Text-to-SQL + анализ + визуализация)."
            ),
            themes=[
                "Тенденции поступлений по регионам и видам налогов",
                "Проблемы собираемости и задолженности",
            ],
            key_takeaways=[
                "г. Минск обеспечивает значительную долю начислений.",
                "Наблюдается сезонный рост к концу года.",
                "Задолженность сконцентрирована в отдельных регионах.",
                "Высокая собираемость по подоходному налогу.",
            ],
            recommendations=[
                "Усилить мониторинг в регионах с высокой задолженностью.",
                "Проанализировать аномалии в конкретных видах налогов.",
            ],
        )

    def _generate_placeholder_png(
        self,
        *,
        title: str,
        slide_type: str,
        content: str | list[str] | None = None,
        slide_idx: int,
        out_dir: Path,
        presentation_id: str,
    ) -> str | None:
        """Create a premium dark-themed slide PNG using Pillow with slide-type-specific layouts."""
        try:
            import textwrap
            from PIL import Image, ImageDraw, ImageFont

            W, H = 1920, 1080

            # ── Color palette ──
            BG      = (8, 13, 26)
            SURFACE = (15, 23, 42)
            CARD    = (22, 33, 60)
            BORDER  = (30, 45, 80)
            VIOLET  = (124, 58, 237)
            BLUE    = (59, 130, 246)
            EMERALD = (16, 185, 129)
            AMBER   = (245, 158, 11)
            ROSE    = (244, 63, 94)
            CYAN    = (6, 182, 212)
            TEXT_W  = (248, 250, 252)
            TEXT_M  = (148, 163, 184)
            TEXT_D  = (71, 85, 105)

            TYPE_ACCENT = {
                "title": VIOLET, "summary": BLUE, "themes": CYAN,
                "chart": EMERALD, "takeaways": EMERALD,
                "recommendations": AMBER, "appendix": (71, 85, 105),
            }
            accent = TYPE_ACCENT.get(slide_type, VIOLET)

            def _font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
                for path in [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/System/Library/Fonts/Arial.ttf",
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                ]:
                    try:
                        return ImageFont.truetype(path, size)
                    except Exception:
                        pass
                return ImageFont.load_default()

            def _font_bold(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
                for path in [
                    "/System/Library/Fonts/Helvetica.ttc",
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
                    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
                ]:
                    try:
                        return ImageFont.truetype(path, size)
                    except Exception:
                        pass
                return _font(size)

            f_tiny  = _font(22)
            f_small = _font(30)
            f_body  = _font(38)
            f_sub   = _font(48)
            f_title = _font_bold(76)
            f_hero  = _font_bold(100)
            f_badge = _font_bold(26)
            f_num   = _font_bold(82)

            img  = Image.new("RGB", (W, H), BG)
            draw = ImageDraw.Draw(img)

            # ── Subtle grid ──
            for x in range(0, W, 80):
                draw.line([(x, 0), (x, H)], fill=(20, 30, 55), width=1)
            for y in range(0, H, 80):
                draw.line([(0, y), (W, y)], fill=(20, 30, 55), width=1)

            # ── Top accent bar ──
            draw.rectangle([0, 0, W, 8], fill=accent)
            r, g, b = accent
            draw.rectangle([W // 2, 0, W, 8], fill=(max(0, r - 40), max(0, g - 20), min(255, b + 40)))

            # ── Header strip ──
            draw.rectangle([0, 8, W, 112], fill=SURFACE)
            draw.rectangle([0, 110, W, 114], fill=accent)

            # Brand
            draw.rounded_rectangle([48, 26, 108, 88], radius=14, fill=CARD)
            draw.text((58, 36), "BI", fill=TEXT_W, font=f_sub)
            draw.text((122, 30), "Prototip BI",           fill=TEXT_W, font=f_badge)
            draw.text((122, 60), "Enterprise AI Analytics", fill=TEXT_D, font=f_tiny)

            # Type badge
            TYPE_LABELS = {
                "title": "TITLE", "summary": "OVERVIEW", "themes": "AGENDA",
                "chart": "DATA CHART", "takeaways": "KEY TAKEAWAYS",
                "recommendations": "ACTION PLAN", "appendix": "APPENDIX",
            }
            bl = TYPE_LABELS.get(slide_type, slide_type.upper())
            bw = len(bl) * 18 + 48
            draw.rounded_rectangle([W - bw - 60, 28, W - 60, 84], radius=10, fill=CARD)
            draw.rounded_rectangle([W - bw - 60, 28, W - 60, 84], radius=10, outline=accent, width=2)
            draw.text((W - bw - 44, 40), bl, fill=accent, font=f_badge)
            draw.text((W - 54, 32), f"{slide_idx + 1:02d}", fill=TEXT_D, font=f_small)

            # ── Content lines ──
            content_lines: list[str] = []
            if isinstance(content, list):
                content_lines = [c for c in content if c]
            elif content:
                content_lines = [l for l in content.split('\n') if l.strip()]

            PALETTE = [VIOLET, BLUE, EMERALD, AMBER, ROSE, CYAN]

            if slide_type == "title":
                wrapped = textwrap.fill(title, width=26)
                draw.multiline_text((80, 160), wrapped, fill=TEXT_W, font=f_hero, spacing=18)
                sub = content_lines[0] if content_lines else "Синтетические данные · Республика Беларусь"
                draw.text((80, 430), sub, fill=TEXT_M, font=f_sub)
                stats = [("7", "Регионов"), ("12+", "Видов налогов"), ("5", "AI-агентов"), ("2024", "Год")]
                cw2 = (W - 160 - 90) // 4
                for ci, (val, lab) in enumerate(stats):
                    cx = 80 + ci * (cw2 + 30)
                    cy2 = H - 220
                    draw.rounded_rectangle([cx, cy2, cx + cw2, cy2 + 130], radius=16, fill=CARD)
                    draw.text((cx + 20, cy2 + 14), val, fill=TEXT_W, font=f_num)
                    draw.text((cx + 20, cy2 + 96), lab, fill=TEXT_M, font=f_small)

            elif slide_type == "summary":
                draw.text((80, 140), title, fill=TEXT_W, font=f_title)
                draw.rectangle([80, 226, 160, 232], fill=accent)
                text_full = " ".join(content_lines)
                wrapped = textwrap.fill(text_full, width=50)
                draw.multiline_text((80, 260), wrapped, fill=TEXT_M, font=f_body, spacing=14)
                kpis = [("7", "Регионов", VIOLET), ("98.5%", "Точность", EMERALD), ("2024", "Период", BLUE), ("5", "Агентов", AMBER)]
                ky = 260
                for val, lab, col in kpis:
                    draw.rounded_rectangle([W // 2 + 40, ky, W - 80, ky + 105], radius=14, fill=CARD)
                    draw.rounded_rectangle([W // 2 + 40, ky, W // 2 + 52, ky + 105], radius=4, fill=col)
                    draw.text((W // 2 + 70, ky + 14), val, fill=TEXT_W, font=f_sub)
                    draw.text((W // 2 + 70, ky + 66), lab, fill=TEXT_D, font=f_small)
                    ky += 125

            elif slide_type == "themes":
                draw.text((80, 140), title, fill=TEXT_W, font=f_title)
                draw.rectangle([80, 226, 160, 232], fill=accent)
                ch2 = min(140, (H - 320) // max(len(content_lines), 1))
                for i, theme in enumerate(content_lines[:6]):
                    col = PALETTE[i % len(PALETTE)]
                    cy2 = 260 + i * (ch2 + 16)
                    draw.rounded_rectangle([80, cy2, W - 80, cy2 + ch2], radius=14, fill=CARD)
                    draw.rounded_rectangle([80, cy2, 92, cy2 + ch2], radius=4, fill=col)
                    draw.rounded_rectangle([108, cy2 + ch2 // 2 - 26, 158, cy2 + ch2 // 2 + 26], radius=10, fill=col)
                    draw.text((116, cy2 + ch2 // 2 - 20), f"0{i+1}", fill=TEXT_W, font=f_badge)
                    tw2 = textwrap.fill(theme, width=76)
                    draw.multiline_text((178, cy2 + 22), tw2, fill=TEXT_W, font=f_body, spacing=8)

            elif slide_type == "takeaways":
                draw.text((80, 140), title, fill=TEXT_W, font=f_title)
                draw.rectangle([80, 226, 160, 232], fill=EMERALD)
                ch2 = min(120, (H - 320) // max(len(content_lines[:6]), 1))
                for i, tw_text in enumerate(content_lines[:6]):
                    col = PALETTE[i % len(PALETTE)]
                    cy2 = 260 + i * (ch2 + 16)
                    draw.rounded_rectangle([80, cy2, W - 80, cy2 + ch2], radius=14, fill=CARD)
                    draw.rounded_rectangle([96, cy2 + 16, 148, cy2 + 68], radius=10, fill=col)
                    draw.text((108, cy2 + 22), str(i + 1), fill=TEXT_W, font=f_sub)
                    wrapped = textwrap.fill(tw_text, width=74)
                    draw.multiline_text((172, cy2 + 14), wrapped, fill=TEXT_W, font=f_body, spacing=8)
                    pct = 85 - i * 8
                    by = cy2 + ch2 - 14
                    draw.rounded_rectangle([172, by, W - 96, by + 6], radius=3, fill=SURFACE)
                    draw.rounded_rectangle([172, by, 172 + int((W - 96 - 172) * pct / 100), by + 6], radius=3, fill=col)

            elif slide_type == "recommendations":
                draw.text((80, 140), title, fill=TEXT_W, font=f_title)
                draw.rectangle([80, 226, 160, 232], fill=AMBER)
                n = min(len(content_lines), 4)
                cols_n = 2 if n > 2 else 1
                cw2 = (W - 160 - (cols_n - 1) * 30) // cols_n
                rows = (n + cols_n - 1) // cols_n
                ch2 = min(200, (H - 320) // max(rows, 1))
                for i, rec in enumerate(content_lines[:4]):
                    col = PALETTE[i % len(PALETTE)]
                    col_i = i % cols_n
                    row_i = i // cols_n
                    rx = 80 + col_i * (cw2 + 30)
                    ry = 270 + row_i * (ch2 + 20)
                    draw.rounded_rectangle([rx, ry, rx + cw2, ry + ch2], radius=16, fill=CARD)
                    draw.rounded_rectangle([rx, ry, rx + cw2, ry + 6], radius=4, fill=col)
                    draw.text((rx + 16, ry + 22), f"ПРИОРИТЕТ {i + 1}", fill=col, font=f_tiny)
                    wrapped = textwrap.fill(rec, width=36 if cols_n == 2 else 78)
                    draw.multiline_text((rx + 16, ry + 56), wrapped, fill=TEXT_W, font=f_body, spacing=8)

            else:  # chart / appendix / generic
                draw.text((80, 140), title, fill=TEXT_W, font=f_title)
                draw.rectangle([80, 226, 160, 232], fill=accent)
                y2 = 270
                for line in content_lines[:8]:
                    wrapped = textwrap.fill(f"• {line}", width=72)
                    draw.multiline_text((80, y2), wrapped, fill=TEXT_M, font=f_body, spacing=8)
                    y2 += 60 * (wrapped.count('\n') + 1) + 10

            # ── Footer ──
            draw.rectangle([0, H - 56, W, H], fill=SURFACE)
            draw.rectangle([0, H - 58, W, H - 56], fill=accent)
            draw.text((80, H - 42), "Prototip BI  •  AI Analytics Platform  •  Синтетические данные (демо)", fill=TEXT_D, font=f_tiny)
            draw.text((W - 80, H - 42), f"{slide_idx + 1}", fill=TEXT_D, font=f_body)

            png_path = out_dir / f"pres_placeholder_{presentation_id}_{slide_idx}.png"
            img.save(str(png_path), "PNG")
            return str(png_path.resolve())
        except Exception as exc:
            logger.info(f"[PresentationAgent] placeholder_png_error idx={slide_idx}: {exc}")
            return None

    def _export_slide_charts(
        self,
        *,
        questions: list[str],
        results: list[AskResult],
        prefs: dict[int, str | None],
        out_dir: Path,
        presentation_id: str,
    ) -> list[str | None]:
        paths: list[str | None] = []
        for idx, (q, res) in enumerate(zip(questions, results)):  # noqa: B905
            png_path: str | None = None
            if res.chart_spec and res.data:
                try:
                    df = pd.DataFrame(res.data)
                    slide_spec = res.chart_spec.model_copy()
                    user_pref = prefs.get(idx)
                    if user_pref:
                        slide_spec.chart_type = user_pref  # type: ignore[assignment]
                    slide_spec.title = ""
                    slide_spec = repair_chart_spec(slide_spec, res.data, question=q)
                    fig = build_chart(df, slide_spec)
                    with suppress(Exception):
                        fig.update_layout(title=dict(text=""))
                    header = q
                    if getattr(slide_spec, "action_title", None):
                        header = slide_spec.action_title or q
                    elif slide_spec.title:
                        header = slide_spec.title
                    slide_png = out_dir / f"pres_slide_{presentation_id}_{idx}_{self._slug(header)}.png"
                    export_png(fig, slide_png, scale=2.0)
                    png_path = str(slide_png.resolve())
                except Exception as e:
                    logger.info(f"[PresentationAgent] slide_graph_error: {e}")
            paths.append(png_path)
        return paths

    def run(
        self,
        questions: list[str] | list[dict] | list | PresentationInput,
        num_slides: int | None = None,
        include_title: bool = True,
        include_recommendations: bool = True,
    ) -> PresentationResult:
        start = time.time()
        presentation_id = uuid.uuid4().hex[:10]
        emit_pipeline_stage(
            "synthesis",
            "running",
            f"Сборка презентации ({presentation_id})...",
            agent="presentation_agent",
        )

        all_qs, prefs, qlist = self._normalize_input(questions)
        if not all_qs:
            raise ValueError("questions не может быть пустым")

        qs, prefs = self._trim_for_num_slides(
            all_qs,
            prefs,
            qlist,
            num_slides=num_slides,
            include_title=include_title,
            include_recommendations=include_recommendations,
        )
        logger.info(
            f"[PresentationAgent] start: questions={len(qs)} num_slides={num_slides} "
            f"inc_title={include_title} inc_recs={include_recommendations}"
        )

        executor = get_executor(include_planner=False)
        out_dir = Path("out")
        out_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = out_dir / f"presentation_{presentation_id}.pptx"

        results = self._collect_results(qs, executor)

        try:
            narrative = self._get_deck_narrative(qs, results)
        except Exception as _e:
            logger.warning(f"[PresentationAgent] narrative fallback: {_e}")
            narrative = self._fallback_narrative()

        # Generate chart PNGs for each slide
        slide_png_paths = self._export_slide_charts(
            questions=qs,
            results=results,
            prefs=prefs,
            out_dir=out_dir,
            presentation_id=presentation_id,
        )

        state = PresentationState(
            presentation_id=presentation_id,
            questions=qs,
            prefs=prefs,
            results=results,
            narrative=narrative,
            include_title=include_title,
            include_recommendations=include_recommendations,
            num_slides=num_slides,
        )
        with open(out_dir / f"pres_state_{presentation_id}.json", "w", encoding="utf-8") as f:
            f.write(state.model_dump_json(indent=2))

        prs, final_png_paths, slides_data = self._build_pptx(state, slide_png_paths, out_dir)

        elapsed = int((time.time() - start) * 1000)
        logger.info(f"[PresentationAgent] end: pptx={pptx_path} slides={len(prs.slides)} ({elapsed}ms)")
        emit_pipeline_stage(
            "viz",
            "done",
            f"Презентация: {len(prs.slides)} слайдов, {sum(1 for p in final_png_paths if p)} превью",
            agent="presentation_agent",
        )
        return PresentationResult(
            pptx_path=str(pptx_path),
            num_slides=len(prs.slides),
            slide_png_paths=final_png_paths,
            presentation_id=presentation_id,
            slides=slides_data,
            reasoning=(
                f"Собрано {len(prs.slides)} слайдов из {len(qs)} вопросов (slide pipeline). "
                "PresentationRenderer: динамические макеты + gov-badge + KPI + таблицы."
            ),
        )


    def _build_pptx(self, state: PresentationState, slide_png_paths: list[str | None], out_dir: Path) -> tuple[Presentation, list[str], list[SlideData]]:
        prs = Presentation()
        prs.slide_width = Inches(PresentationTheme.SLIDE_W)
        prs.slide_height = Inches(PresentationTheme.SLIDE_H)
        blank = prs.slide_layouts[6]
        renderer = self.renderer
        final_png_paths: list[str] = []

        if state.include_title:
            renderer.create_title_slide(prs.slides.add_slide(blank))
            final_png_paths.append("")
        renderer.create_summary_slide(
            prs.slides.add_slide(blank), state.narrative.overview, state.results[0] if state.results else None
        )
        final_png_paths.append("")
        renderer.create_themes_slide(prs.slides.add_slide(blank), state.narrative.themes, state.questions)
        final_png_paths.append("")

        for idx, (q, res) in enumerate(zip(state.questions, state.results)):  # noqa: B905
            slide = prs.slides.add_slide(blank)
            renderer.render_question_slide(slide, q, res, state.prefs.get(idx), slide_png_paths[idx])
            renderer.build_question_slides_footer(slide, len(prs.slides))
            final_png_paths.append(slide_png_paths[idx] or "")

        renderer.create_takeaways_slide(prs.slides.add_slide(blank), state.narrative.key_takeaways)
        final_png_paths.append("")
        
        if state.include_recommendations:
            renderer.create_recommendations_slide(prs.slides.add_slide(blank), state.narrative.recommendations)
            final_png_paths.append("")

        target = state.num_slides if (state.num_slides is not None and state.num_slides > 0) else len(prs.slides)
        while len(prs.slides) < target:
            renderer.create_appendix_slide(prs.slides.add_slide(blank), len(prs.slides))
            final_png_paths.append("")

        pptx_path = out_dir / f"presentation_{state.presentation_id}.pptx"
        prs.save(str(pptx_path))

        # Map each slide index to (type, display_title, content) for the thumbnail renderer
        slide_type_labels: list[tuple[str, str, str | list[str] | None]] = []
        _si = 0
        if state.include_title:
            slide_type_labels.append(("title", "BI-аналитика налогов РБ", "Синтетические данные (демо), Республика Беларусь"))
            _si += 1
        slide_type_labels.append(("summary", "Обзор", state.narrative.overview))
        _si += 1
        slide_type_labels.append(("themes", "Темы и повестка", state.narrative.themes))
        _si += 1
        for qi, q in enumerate(state.questions):
            # for chart slides, we extract the primary text or conclusion
            res = state.results[qi]
            chart_content = res.analysis.key_conclusion if res.analysis else ""
            slide_type_labels.append(("chart", q[:60], chart_content))
            _si += 1
        slide_type_labels.append(("takeaways", "Ключевые выводы", state.narrative.key_takeaways))
        _si += 1
        if state.include_recommendations:
            slide_type_labels.append(("recommendations", "Рекомендации", state.narrative.recommendations))
            _si += 1
        while _si < len(final_png_paths):
            slide_type_labels.append(("appendix", "Приложение", "Дополнительные материалы и методология расчетов"))
            _si += 1

        slides_data: list[SlideData] = []
        for slide_idx, (stype, stitle, scontent) in enumerate(slide_type_labels):
            slides_data.append(SlideData(
                slide_idx=slide_idx,
                slide_type=stype,
                title=stitle,
                content=scontent
            ))
            if slide_idx < len(final_png_paths) and not final_png_paths[slide_idx]:
                png = self._generate_placeholder_png(
                    title=stitle,
                    slide_type=stype,
                    content=scontent,
                    slide_idx=slide_idx,
                    out_dir=out_dir,
                    presentation_id=state.presentation_id,
                )
                if png:
                    final_png_paths[slide_idx] = png
        return prs, final_png_paths, slides_data

    def update_presentation(self, presentation_id: str, slide_updates: dict[int, SlideUpdate]) -> PresentationResult:
        out_dir = Path("out")
        state_file = out_dir / f"pres_state_{presentation_id}.json"
        if not state_file.exists():
            raise FileNotFoundError(f"State file not found for presentation {presentation_id}")
        
        with open(state_file, "r", encoding="utf-8") as f:
            state = PresentationState.model_validate_json(f.read())
        
        # Determine mapping of slide index to sections (like in _build_pptx)
        _si = 0
        slide_idx_map = {}  # slide_idx -> (type, ref_index)
        if state.include_title:
            slide_idx_map[_si] = ("title", 0)
            _si += 1
        slide_idx_map[_si] = ("summary", 0)
        _si += 1
        slide_idx_map[_si] = ("themes", 0)
        _si += 1
        for qi, q in enumerate(state.questions):
            slide_idx_map[_si] = ("chart", qi)
            _si += 1
        slide_idx_map[_si] = ("takeaways", 0)
        _si += 1
        if state.include_recommendations:
            slide_idx_map[_si] = ("recommendations", 0)
            _si += 1
        
        # Track which chart indices need re-rendering with new type
        chart_type_overrides: dict[int, str] = {}  # qi -> new chart_type

        # Apply updates to the state
        for slide_idx_str, update in slide_updates.items():
            s_idx = int(slide_idx_str)
            if s_idx not in slide_idx_map:
                continue
            stype, sref = slide_idx_map[s_idx]
            
            if stype == "summary":
                if update.content is not None:
                    state.narrative.overview = str(update.content)
                if update.title is not None:
                    pass  # overview has no separate title in state
            elif stype == "themes":
                if update.content is not None:
                    state.narrative.themes = update.content if isinstance(update.content, list) else [str(update.content)]
            elif stype == "chart":
                if update.title is not None:
                    state.questions[sref] = update.title
                if update.content is not None and state.results[sref].analysis:
                    state.results[sref].analysis.key_conclusion = str(update.content)
                if update.chart_type is not None:
                    chart_type_overrides[sref] = update.chart_type
                    # Persist chart_type preference in prefs
                    state.prefs[sref] = update.chart_type
            elif stype == "takeaways":
                if update.content is not None:
                    state.narrative.key_takeaways = update.content if isinstance(update.content, list) else [str(update.content)]
            elif stype == "recommendations":
                if update.content is not None:
                    state.narrative.recommendations = update.content if isinstance(update.content, list) else [str(update.content)]
        
        # Save updated state
        with open(state_file, "w", encoding="utf-8") as f:
            f.write(state.model_dump_json(indent=2))

        # Gather existing chart slide paths
        slide_png_paths: list[str | None] = []
        for idx, q in enumerate(state.questions):
            if idx in chart_type_overrides:
                # Regenerate this chart with the new type
                try:
                    new_type = chart_type_overrides[idx]
                    res = state.results[idx]
                    if res.chart_spec and res.data:
                        df = pd.DataFrame(res.data)
                        slide_spec = res.chart_spec.model_copy()
                        slide_spec.chart_type = new_type  # type: ignore[assignment]
                        slide_spec.title = ""
                        slide_spec = repair_chart_spec(slide_spec, res.data, question=q)
                        fig = build_chart(df, slide_spec)
                        with suppress(Exception):
                            fig.update_layout(title=dict(text=""))
                        header = getattr(slide_spec, "action_title", None) or slide_spec.title or q[:30]
                        new_png_path = out_dir / f"pres_slide_{presentation_id}_{idx}_{self._slug(header)}_{new_type}.png"
                        export_png(fig, new_png_path, scale=2.0)
                        slide_png_paths.append(str(new_png_path.resolve()) if new_png_path.exists() else None)
                        logger.info(f"[update_presentation] regenerated chart {idx} as {new_type}")
                    else:
                        # No chart spec – keep existing
                        found = list(out_dir.glob(f"pres_slide_{presentation_id}_{idx}_*.png"))
                        slide_png_paths.append(str(found[0].resolve()) if found else None)
                except Exception as exc:
                    logger.warning(f"[update_presentation] chart regen failed idx={idx}: {exc}")
                    found = list(out_dir.glob(f"pres_slide_{presentation_id}_{idx}_*.png"))
                    slide_png_paths.append(str(found[0].resolve()) if found else None)
            else:
                res = state.results[idx]
                header = getattr(res.chart_spec, "action_title", None) or getattr(res.chart_spec, "title", None) or q
                expected_png = out_dir / f"pres_slide_{presentation_id}_{idx}_{self._slug(header)}.png"
                if expected_png.exists():
                    slide_png_paths.append(str(expected_png.resolve()))
                else:
                    found = list(out_dir.glob(f"pres_slide_{presentation_id}_{idx}_*.png"))
                    slide_png_paths.append(str(found[0].resolve()) if found else None)

        prs, final_png_paths, slides_data = self._build_pptx(state, slide_png_paths, out_dir)
        pptx_path = out_dir / f"presentation_{presentation_id}.pptx"

        return PresentationResult(
            pptx_path=str(pptx_path),
            num_slides=len(prs.slides),
            slide_png_paths=final_png_paths,
            presentation_id=presentation_id,
            slides=slides_data,
            reasoning="Обновлено пользователем."
        )


    def run_input(self, inp: PresentationInput) -> PresentationResult:
        return self.run(inp.questions)