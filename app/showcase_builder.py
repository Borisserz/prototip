"""Офлайн-сборка портфолио showcase/ (графики + презентации без Ollama)."""

from __future__ import annotations

import json
import re
import tempfile
import uuid
from contextlib import suppress
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from app.agents.models import AnalysisResult, AskResult, DeckNarrative
from app.chart_repair import repair_chart_spec
from app.showcase_catalog import (
    ChartShowcaseEntry,
    PresentationBundle,
    chart_showcase_entries,
    presentation_bundles,
    waterfall_demo_df,
)
from app.presentation_renderer import PresentationRenderer, PresentationTheme
from core.models import ChartSpec
from viz.charts import build_chart, export_html, export_png

CHART_SCALE = 2.5


def _manifest_rel(path: Path, base: Path) -> str:
    """Путь для manifest.json — относительно base (портативно для git)."""
    try:
        return path.resolve().relative_to(base.resolve()).as_posix()
    except ValueError:
        return path.as_posix()


SPEC_OVERRIDE_KEYS = frozenset(
    {
        "title",
        "subtitle",
        "x",
        "y",
        "color",
        "agg",
        "top_n",
        "sort_order",
        "show_average",
        "highlight_category",
        "action_title",
        "insights",
        "filter_tax_type",
        "use_waterfall_df",
    }
)


def load_base_dataset(csv_path: Path | None = None) -> pd.DataFrame:
    path = csv_path or Path("data/sample.csv")
    return pd.read_csv(path)


def _slug(text: str, max_len: int = 40) -> str:
    slug = re.sub(r"[^a-zA-Zа-яА-Я0-9]+", "_", text.lower()).strip("_")
    return slug[:max_len] or "chart"


def catalog_by_chart_type() -> dict[str, ChartShowcaseEntry]:
    return {e.slug: e for e in chart_showcase_entries()}


def prepare_entry_df(entry: ChartShowcaseEntry, base_df: pd.DataFrame) -> pd.DataFrame:
    if entry.use_waterfall_df:
        return waterfall_demo_df()
    df = base_df.copy()
    if entry.filter_tax_type and "tax_type" in df.columns:
        df = df[df["tax_type"] == entry.filter_tax_type]
    return df


def _question_overrides(question: dict[str, Any]) -> dict[str, Any]:
    return {k: question[k] for k in SPEC_OVERRIDE_KEYS if k in question}


def resolve_question_entry(
    question: dict[str, Any], catalog: dict[str, ChartShowcaseEntry]
) -> ChartShowcaseEntry:
    chart_type = str(question.get("chart_type") or "bar")
    if chart_type in catalog:
        return catalog[chart_type]
    for entry in chart_showcase_entries():
        if entry.spec.chart_type == chart_type:
            return entry
    return catalog["bar"]


def build_spec_for_question(
    question: dict[str, Any],
    catalog: dict[str, ChartShowcaseEntry],
) -> tuple[ChartSpec, ChartShowcaseEntry]:
    entry = resolve_question_entry(question, catalog)
    overrides = _question_overrides(question)
    use_waterfall = overrides.pop("use_waterfall_df", entry.use_waterfall_df)
    filter_tax = overrides.pop("filter_tax_type", entry.filter_tax_type)

    spec = entry.spec.model_copy(update=overrides)
    patched_entry = ChartShowcaseEntry(
        index=entry.index,
        slug=entry.slug,
        spec=spec,
        filter_tax_type=filter_tax,
        use_waterfall_df=use_waterfall,
    )
    return spec, patched_entry


def _analysis_from_spec(spec: ChartSpec) -> AnalysisResult:
    insights = list(spec.insights[:3])
    while len(insights) < 3:
        insights.append("Данные синтетические (демо), Республика Беларусь.")
    return AnalysisResult(
        insights=insights,
        key_conclusion=spec.action_title or spec.title,
        anomaly_or_trend="Сезонный рост к концу года; аномалия в Гомельской области (НДС, сентябрь).",
        reasoning="Офлайн-нарратив showcase (без LLM).",
    )


def build_ask_result(
    question_text: str,
    spec: ChartSpec,
    df: pd.DataFrame,
) -> AskResult:
    data = df.to_dict(orient="records")
    repaired = repair_chart_spec(spec.model_copy(), data, question=question_text)
    return AskResult(
        question=question_text,
        sql="-- showcase offline",
        data=data,
        analysis=_analysis_from_spec(repaired),
        chart_spec=repaired,
        reasoning="Showcase offline AskResult",
    )


def export_chart_assets(
    entry: ChartShowcaseEntry,
    df: pd.DataFrame,
    out_dir: Path,
    *,
    scale: float = CHART_SCALE,
    manifest_base: Path | None = None,
) -> dict[str, str]:
    out_dir.mkdir(parents=True, exist_ok=True)
    fig = build_chart(df, entry.spec)
    png_path = out_dir / "chart.png"
    html_path = out_dir / "chart.html"
    spec_path = out_dir / "spec.json"
    export_png(fig, png_path, scale=scale)
    export_html(fig, html_path)
    spec_path.write_text(
        entry.spec.model_dump_json(indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    base = manifest_base or out_dir.parent.parent
    return {
        "png": _manifest_rel(png_path, base),
        "html": _manifest_rel(html_path, base),
        "spec": _manifest_rel(spec_path, base),
    }


def _trim_questions(
    bundle: PresentationBundle,
    *,
    include_title: bool = True,
    include_recommendations: bool = True,
) -> list[dict[str, Any]]:
    base_fixed = (1 if include_title else 0) + 2 + 1 + (1 if include_recommendations else 0)
    max_q = max(0, bundle.num_slides - base_fixed)
    questions = bundle.questions[:max_q] if max_q > 0 else bundle.questions[:1]
    return questions


def _export_slide_pngs(
    questions: list[dict[str, Any]],
    results: list[AskResult],
    out_dir: Path,
    presentation_id: str,
) -> list[str | None]:
    paths: list[str | None] = []
    for idx, (qdict, res) in enumerate(zip(questions, results)):  # noqa: B905
        png_path: str | None = None
        if res.chart_spec and res.data:
            try:
                df = pd.DataFrame(res.data)
                slide_spec = res.chart_spec.model_copy()
                pref = qdict.get("chart_type")
                if pref:
                    slide_spec.chart_type = pref  # type: ignore[assignment]
                slide_spec.title = ""
                slide_spec = repair_chart_spec(slide_spec, res.data, question=str(qdict.get("text", "")))
                fig = build_chart(df, slide_spec)
                with suppress(Exception):
                    fig.update_layout(title=dict(text=""))
                header = str(qdict.get("text", ""))
                if getattr(slide_spec, "action_title", None):
                    header = slide_spec.action_title or header
                elif slide_spec.title:
                    header = slide_spec.title
                slide_png = out_dir / f"pres_slide_{presentation_id}_{idx}_{_slug(header)}.png"
                export_png(fig, slide_png, scale=2.0)
                png_path = str(slide_png.resolve())
            except Exception:
                png_path = None
        paths.append(png_path)
    return paths


def build_offline_presentation(
    bundle: PresentationBundle,
    base_df: pd.DataFrame,
    out_path: Path,
    *,
    include_title: bool = True,
    include_recommendations: bool = True,
) -> dict[str, Any]:
    catalog = catalog_by_chart_type()
    questions = _trim_questions(
        bundle,
        include_title=include_title,
        include_recommendations=include_recommendations,
    )
    results: list[AskResult] = []
    prefs: dict[int, str | None] = {}
    for idx, qdict in enumerate(questions):
        text = str(qdict.get("text", "")).strip()
        spec, entry = build_spec_for_question(qdict, catalog)
        df = prepare_entry_df(entry, base_df)
        results.append(build_ask_result(text, spec, df))
        prefs[idx] = qdict.get("chart_type")

    narrative = DeckNarrative(
        overview=bundle.overview,
        themes=bundle.themes,
        key_takeaways=bundle.key_takeaways,
        recommendations=bundle.recommendations,
        reasoning="Showcase offline DeckNarrative",
    )

    presentation_id = uuid.uuid4().hex[:10]
    qs_text = [str(q.get("text", "")) for q in questions]
    with tempfile.TemporaryDirectory(prefix="showcase_slides_") as slide_tmp:
        slide_dir = Path(slide_tmp)
        slide_png_paths = _export_slide_pngs(questions, results, slide_dir, presentation_id)

        renderer = PresentationRenderer()
        prs = Presentation()
        prs.slide_width = Inches(PresentationTheme.SLIDE_W)
        prs.slide_height = Inches(PresentationTheme.SLIDE_H)
        blank = prs.slide_layouts[6]

        if include_title:
            renderer.create_title_slide(prs.slides.add_slide(blank))
        renderer.create_summary_slide(
            prs.slides.add_slide(blank), narrative.overview, results[0] if results else None
        )
        renderer.create_themes_slide(prs.slides.add_slide(blank), narrative.themes, qs_text)

        for idx, (qdict, res) in enumerate(zip(questions, results)):  # noqa: B905
            slide = prs.slides.add_slide(blank)
            renderer.render_question_slide(
                slide, str(qdict.get("text", "")), res, prefs.get(idx), slide_png_paths[idx]
            )
            renderer.build_question_slides_footer(slide, len(prs.slides))

        renderer.create_takeaways_slide(prs.slides.add_slide(blank), narrative.key_takeaways)
        if include_recommendations:
            renderer.create_recommendations_slide(
                prs.slides.add_slide(blank), narrative.recommendations
            )

        target = bundle.num_slides if bundle.num_slides > 0 else len(prs.slides)
        while len(prs.slides) < target:
            renderer.create_appendix_slide(prs.slides.add_slide(blank), len(prs.slides))

        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        num_slides = len(prs.slides)
        exported_pngs = [p for p in slide_png_paths if p]

    return {
        "filename": bundle.filename,
        "title": bundle.title,
        "pptx": str(out_path.resolve()),
        "num_slides": num_slides,
        "questions": len(questions),
        "slide_pngs": exported_pngs,
    }


def _presentation_manifest_entry(
    meta: dict[str, Any],
    *,
    showcase_root: Path,
) -> dict[str, Any]:
    """Запись презентации для manifest.json (без временных slide_pngs)."""
    pptx = Path(str(meta["pptx"]))
    return {
        "filename": meta["filename"],
        "title": meta["title"],
        "pptx": _manifest_rel(pptx, showcase_root),
        "num_slides": meta["num_slides"],
        "questions": meta["questions"],
    }


def generate_showcase(
    root: Path | str = "showcase",
    *,
    csv_path: Path | None = None,
    chart_scale: float = CHART_SCALE,
) -> dict[str, Any]:
    """Генерирует showcase/charts, showcase/presentations и manifest.json."""
    root_path = Path(root)
    charts_root = root_path / "charts"
    pres_root = root_path / "presentations"
    charts_root.mkdir(parents=True, exist_ok=True)
    pres_root.mkdir(parents=True, exist_ok=True)

    base_df = load_base_dataset(csv_path)
    chart_manifest: list[dict[str, Any]] = []
    for entry in chart_showcase_entries():
        folder = charts_root / f"{entry.index:02d}_{entry.slug}"
        df = prepare_entry_df(entry, base_df)
        assets = export_chart_assets(
            entry, df, folder, scale=chart_scale, manifest_base=root_path
        )
        chart_manifest.append(
            {
                "index": entry.index,
                "slug": entry.slug,
                "chart_type": entry.spec.chart_type,
                "title": entry.spec.title,
                "action_title": entry.spec.action_title,
                "folder": str(folder.relative_to(root_path)),
                **assets,
            }
        )

    presentation_manifest: list[dict[str, Any]] = []
    for bundle in presentation_bundles():
        out_pptx = pres_root / bundle.filename
        meta = build_offline_presentation(bundle, base_df, out_pptx)
        presentation_manifest.append(_presentation_manifest_entry(meta, showcase_root=root_path))

    dataset_path = csv_path or Path("data/sample.csv")
    repo_root = root_path.resolve().parent
    try:
        dataset_ref = dataset_path.resolve().relative_to(repo_root).as_posix()
    except ValueError:
        dataset_ref = dataset_path.as_posix()

    manifest = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source": "scripts/generate_leadership_showcase.py",
        "dataset": dataset_ref,
        "chart_scale": chart_scale,
        "charts": chart_manifest,
        "presentations": presentation_manifest,
    }
    manifest_path = root_path / "manifest.json"
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    manifest["manifest"] = _manifest_rel(manifest_path, root_path)
    return manifest