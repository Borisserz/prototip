"""Программная верстка презентаций (python-pptx), pixel-perfect к веб-интерфейсу."""

from __future__ import annotations

import enum
from datetime import datetime
from pathlib import Path
from typing import Any

import pandas as pd
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.agents.models import AskResult
from app.kpi_utils import compute_overview_kpis
from core.models import ChartSpec
from viz.style import format_number_ru, get_russian_label

FOOTER_TEXT = "Prototip BI  •  AI Analytics Platform  •  Министерство по налогам и сборам РБ"
BADGE_TEXT = "ДЛЯ СЛУЖЕБНОГО ПОЛЬЗОВАНИЯ | AI DRAFT"

CHART_TYPE_RU: dict[str, str] = {
    "bar": "столбчатая",
    "grouped_bar": "группированная столбчатая",
    "stacked_bar": "стековая столбчатая",
    "line": "линейная",
    "horizontal_bar": "горизонтальная столбчатая",
    "donut": "круговая",
    "kpi": "KPI-индикатор",
    "heatmap": "тепловая карта",
    "area": "областная",
    "scatter": "точечная",
    "waterfall": "водопад",
    "treemap": "древовидная (treemap)",
}


class LayoutKind(str, enum.Enum):
    TOP_BOTTOM = "top_bottom"
    SIDE_BY_SIDE = "side_by_side"
    CHART_TABLE = "chart_table"
    KPI_FULL = "kpi_full"


class T:
    """Design tokens — mirror of the React webapp colours."""
    # Backgrounds
    BG_DARK    = RGBColor(8, 13, 26)         # #080d1a  (webapp exact)
    BG_DARKER  = RGBColor(2, 6, 23)
    CARD_BG    = RGBColor(15, 23, 42)        # slate-900
    CARD_BG2   = RGBColor(20, 30, 55)
    CARD_DARK  = RGBColor(10, 15, 35)
    FOOTER_BG  = RGBColor(6, 10, 20)
    BORDER     = RGBColor(30, 40, 65)

    # Accent colours
    VIOLET     = RGBColor(124, 58, 237)
    VIOLET_L   = RGBColor(167, 139, 250)
    VIOLET_BG  = RGBColor(30, 18, 60)
    VIOLET_B   = RGBColor(60, 35, 100)

    BLUE       = RGBColor(59, 130, 246)
    BLUE_L     = RGBColor(147, 197, 253)
    BLUE_BG    = RGBColor(18, 35, 65)
    BLUE_B     = RGBColor(35, 65, 110)

    EMERALD    = RGBColor(16, 185, 129)
    EMERALD_L  = RGBColor(110, 231, 183)
    EMERALD_BG = RGBColor(10, 40, 30)
    EMERALD_B  = RGBColor(20, 70, 55)

    AMBER      = RGBColor(245, 158, 11)
    AMBER_L    = RGBColor(252, 211, 77)
    AMBER_BG   = RGBColor(45, 30, 10)
    AMBER_B    = RGBColor(80, 55, 20)

    ROSE       = RGBColor(244, 63, 94)
    ROSE_BG    = RGBColor(45, 10, 20)
    ROSE_B     = RGBColor(80, 25, 40)

    CYAN       = RGBColor(6, 182, 212)
    CYAN_BG    = RGBColor(8, 35, 45)
    CYAN_B     = RGBColor(15, 65, 80)

    # Text
    WHITE          = RGBColor(255, 255, 255)
    TEXT_PRIMARY   = RGBColor(248, 250, 252)
    TEXT_SECONDARY = RGBColor(148, 163, 184)
    TEXT_MUTED     = RGBColor(100, 116, 139)
    TEXT_DIMMED    = RGBColor(71, 85, 105)

    # Table
    TABLE_HEADER_BG = RGBColor(124, 58, 237)
    TABLE_ROW_ALT   = RGBColor(20, 30, 55)
    TABLE_ROW_BASE  = RGBColor(8, 13, 26)

    # Geometry
    SLIDE_W   = 13.333
    SLIDE_H   = 7.5
    M         = 0.45        # margin horizontal
    USABLE_W  = 12.433
    CONTENT_TOP = 1.10
    CONTENT_H   = 5.50
    FOOTER_TOP  = 7.06
    FOOTER_H    = 0.44
    GAP       = 0.12

    FONT = "Arial"

    # Accent colour palettes (indexed 0-5, cycles in loops)
    ACCENTS = [
        {"fg": RGBColor(124, 58, 237),  "bg": RGBColor(30, 18, 60),  "border": RGBColor(60, 35, 100),  "light": RGBColor(167, 139, 250)},  # violet
        {"fg": RGBColor(59, 130, 246),  "bg": RGBColor(18, 35, 65),  "border": RGBColor(35, 65, 110),  "light": RGBColor(147, 197, 253)},  # blue
        {"fg": RGBColor(16, 185, 129),  "bg": RGBColor(10, 40, 30),  "border": RGBColor(20, 70, 55),   "light": RGBColor(110, 231, 183)},  # emerald
        {"fg": RGBColor(245, 158, 11),  "bg": RGBColor(45, 30, 10),  "border": RGBColor(80, 55, 20),   "light": RGBColor(252, 211, 77)},   # amber
        {"fg": RGBColor(244, 63, 94),   "bg": RGBColor(45, 10, 20),  "border": RGBColor(80, 25, 40),   "light": RGBColor(253, 164, 175)},  # rose
        {"fg": RGBColor(6, 182, 212),   "bg": RGBColor(8, 35, 45),   "border": RGBColor(15, 65, 80),   "light": RGBColor(103, 232, 249)},  # cyan
    ]

    # Unicode icon symbols — used instead of SVG/Lucide icons
    ICONS = {
        "pin":       "◉",   # MapPin
        "cpu":       "⊞",   # Cpu
        "target":    "◎",   # Target
        "activity":  "↻",   # Activity
        "trend_up":  "↑",   # TrendingUp
        "trend_dn":  "↓",   # TrendingDown
        "alert":     "⚠",   # AlertTriangle
        "check":     "✓",   # CheckCircle
        "bulb":      "✦",   # Lightbulb
        "arrow":     "→",   # ArrowRight
        "zap":       "★",   # Zap
        "shield":    "◈",   # Shield
        "globe":     "⊙",   # Globe
        "file":      "≡",   # FileText
        "bar_chart": "▌",   # BarChart2
        "user":      "◑",   # Users
        "dollar":    "◆",   # DollarSign
        "star":      "★",
        "award":     "◈",
    }


# Aliases for backward compat
PresentationTheme = T


class PresentationRenderer:
    """Единая точка программной отрисовки слайдов — pixel-perfect к веб-приложению."""

    def __init__(self, *, theme: type[T] = T) -> None:
        self.t = theme

    # ──────────────────────────────────────────────────────────── utils ──────
    @staticmethod
    def choose_layout(chart_type: str | None) -> LayoutKind:
        ct = chart_type or "bar"
        if ct in {"line", "area", "waterfall", "heatmap"}:
            return LayoutKind.TOP_BOTTOM
        if ct in {"donut", "treemap"}:
            return LayoutKind.SIDE_BY_SIDE
        if ct in {"bar", "horizontal_bar", "grouped_bar", "stacked_bar", "scatter"}:
            return LayoutKind.CHART_TABLE
        if ct == "kpi":
            return LayoutKind.KPI_FULL
        return LayoutKind.TOP_BOTTOM

    @staticmethod
    def truncate_text(text: str, max_chars: int = 140) -> str:
        text = (text or "").strip()
        return text if len(text) <= max_chars else text[: max_chars - 3] + "..."

    def _bullet_font_size(self, texts: list[str], base: int = 14) -> int:
        total = sum(len(t) for t in texts)
        if total > 600: return 9
        if total > 450: return 10
        if total > 300: return 11
        if total > 180: return 12
        return base

    def _resolve_headers(self, question: str, res: AskResult) -> tuple[str, str | None]:
        spec = res.chart_spec
        if spec and getattr(spec, "action_title", None):
            return self.truncate_text(spec.action_title, 80), question
        if spec and getattr(spec, "title", None):
            return self.truncate_text(spec.title, 80), question
        return self.truncate_text(question, 80), None

    # ──────────────────────────────────────────────────── low-level draw ─────
    def _bg(self, slide: Any) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.t.BG_DARK

    def _rect(
        self, slide: Any, left: float, top: float, width: float, height: float,
        *, fill: RGBColor | None = None, border: RGBColor | None = None,
        border_pt: float = 0.75, no_line: bool = False,
    ) -> Any:
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        else:
            shp.fill.background()
        if no_line or border is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = border
            shp.line.width = Pt(border_pt)
        return shp

    def _textbox(
        self, slide: Any, left: float, top: float, width: float, height: float,
        text: str = "", *, font_size: int = 12, bold: bool = False,
        color: RGBColor | None = None, align: Any = PP_ALIGN.LEFT,
        wrap: bool = True, italic: bool = False,
    ) -> Any:
        tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tx.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        if text:
            p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = self.t.FONT
        p.font.color.rgb = color or self.t.TEXT_PRIMARY
        return tx

    def _icon_box(
        self, slide: Any, left: float, top: float, size: float,
        symbol: str, fill: RGBColor, color: RGBColor | None = None,
    ) -> None:
        """Draw a colored square with a Unicode symbol — replaces Lucide icons."""
        self._rect(slide, left, top, size, size, fill=fill, no_line=True)
        self._textbox(
            slide, left, top, size, size,
            symbol,
            font_size=int(size * 55),   # scale roughly with size in inches
            bold=True,
            color=color or self.t.WHITE,
            align=PP_ALIGN.CENTER,
        )

    def _gradient_bar(
        self, slide: Any, top: float, height: float = 0.07,
        colors: list[RGBColor] | None = None,
    ) -> None:
        """Gradient bar simulated with N colour segments."""
        cols = colors or [
            self.t.VIOLET, RGBColor(90, 90, 240), self.t.BLUE,
            RGBColor(30, 160, 200), self.t.EMERALD,
        ]
        seg_w = self.t.SLIDE_W / len(cols)
        for i, c in enumerate(cols):
            self._rect(slide, i * seg_w, top, seg_w, height, fill=c, no_line=True)

    def _footer_bar(self, slide: Any, slide_num: int | None = None) -> None:
        self._rect(slide, 0, self.t.FOOTER_TOP, self.t.SLIDE_W, self.t.FOOTER_H,
                   fill=self.t.FOOTER_BG, no_line=True)
        self._rect(slide, 0, self.t.FOOTER_TOP, self.t.SLIDE_W, 0.005,
                   fill=self.t.BORDER, no_line=True)
        self._textbox(
            slide, self.t.M, self.t.FOOTER_TOP + 0.09,
            10.0, 0.28, FOOTER_TEXT,
            font_size=8, color=self.t.TEXT_DIMMED,
        )
        if slide_num is not None:
            self._textbox(
                slide, 11.0, self.t.FOOTER_TOP + 0.09,
                2.0, 0.28, f"Слайд {slide_num}",
                font_size=8, color=self.t.TEXT_DIMMED, align=PP_ALIGN.RIGHT,
            )

    def _gov_badge(self, slide: Any) -> None:
        self._rect(slide, 10.55, 0.12, 2.35, 0.30, fill=self.t.CARD_BG, border=self.t.BORDER)
        self._textbox(slide, 10.57, 0.14, 2.30, 0.26, BADGE_TEXT,
                      font_size=6.5, color=self.t.TEXT_MUTED, align=PP_ALIGN.CENTER)

    def _badge_tag(
        self, slide: Any, left: float, top: float, text: str,
        fill: RGBColor | None = None, border: RGBColor | None = None,
        color: RGBColor | None = None,
    ) -> float:
        """Draw a badge tag; returns its width."""
        fill = fill or self.t.BLUE_BG
        border = border or self.t.BLUE_B
        color = color or self.t.BLUE_L
        width = max(0.80, len(text) * 0.095 + 0.28)
        self._rect(slide, left, top, width, 0.28, fill=fill, border=border, border_pt=0.5)
        self._textbox(slide, left + 0.06, top + 0.04, width - 0.12, 0.22, text,
                      font_size=7.5, bold=True, color=color, align=PP_ALIGN.CENTER)
        return width

    def _card(
        self, slide: Any, left: float, top: float, width: float, height: float,
        fill: RGBColor | None = None, border: RGBColor | None = None,
    ) -> None:
        self._rect(slide, left, top, width, height,
                   fill=fill or self.t.CARD_BG, border=border or self.t.BORDER)

    def _progress_bar(
        self, slide: Any, left: float, top: float, width: float,
        pct: float, accent: dict | None = None,
    ) -> None:
        ac = accent or self.t.ACCENTS[0]
        self._rect(slide, left, top, width, 0.048, fill=self.t.CARD_BG2, no_line=True)
        if pct > 0:
            self._rect(slide, left, top, width * min(pct, 1.0), 0.048, fill=ac["fg"], no_line=True)

    def _place_chart_image(
        self, slide: Any, png_path: str | Path,
        *, left: float, top: float, width: float, height: float,
    ) -> bool:
        """Place chart PNG on slide. Returns True if successful."""
        path = Path(png_path) if png_path else None
        if path and path.exists() and path.stat().st_size > 0:
            slide.shapes.add_picture(
                str(path), Inches(left), Inches(top),
                width=Inches(width), height=Inches(height),
            )
            return True
        return False

    # ──────────────────────────────────────────────────── TABLE helpers ──────
    def _style_table_cell(
        self, cell: Any, text: str, *,
        bold: bool = False, font_size: int = 11,
        fill: RGBColor | None = None, font_color: RGBColor | None = None,
    ) -> None:
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill or self.t.TABLE_ROW_BASE
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(font_size)
            para.font.bold = bold
            para.font.name = self.t.FONT
            para.font.color.rgb = font_color or self.t.TEXT_PRIMARY

    def _build_ranking_table_rows(self, res: AskResult, spec: ChartSpec) -> list[tuple[str, str]]:
        if not res.data:
            return []
        try:
            df = pd.DataFrame(res.data)
            x_col = spec.x
            y_col = spec.y if isinstance(spec.y, str) else (spec.y[0] if spec.y else None)
            if not x_col or not y_col or x_col not in df.columns or y_col not in df.columns:
                return []
            agg = spec.agg or "sum"
            g = (df.groupby(x_col, as_index=False)[y_col].mean()
                 if agg == "mean" else df.groupby(x_col, as_index=False)[y_col].sum())
            g = g.sort_values(y_col, ascending=False).head(5)
            suffix = "бел. руб." if any(k in y_col.lower() for k in ("debt", "accrued", "paid", "penalt")) else ""
            return [(str(row[x_col]), format_number_ru(float(row[y_col]), suffix=suffix))
                    for _, row in g.iterrows()]
        except Exception:
            return []

    def _render_ranking_table(
        self, slide: Any, res: AskResult, spec: ChartSpec,
        *, left: float, top: float, width: float, height: float = 2.70,
    ) -> None:
        rows_data = self._build_ranking_table_rows(res, spec)
        if not rows_data:
            return
        n_rows = min(6, len(rows_data) + 1)
        tbl = slide.shapes.add_table(
            n_rows, 2, Inches(left), Inches(top), Inches(width), Inches(height),
        ).table
        tbl.columns[0].width = Inches(width * 0.55)
        tbl.columns[1].width = Inches(width * 0.45)
        y_col = spec.y if isinstance(spec.y, str) else (spec.y[0] if spec.y else "")
        self._style_table_cell(tbl.cell(0, 0), get_russian_label(spec.x).split(",")[0],
                               bold=True, fill=self.t.TABLE_HEADER_BG, font_color=self.t.WHITE)
        self._style_table_cell(tbl.cell(0, 1), f"{get_russian_label(y_col).split(',')[0]}, бел. руб.",
                               bold=True, fill=self.t.TABLE_HEADER_BG, font_color=self.t.WHITE)
        for i, (cat, val) in enumerate(rows_data[:5], start=1):
            fill = self.t.TABLE_ROW_ALT if i % 2 == 0 else self.t.TABLE_ROW_BASE
            self._style_table_cell(tbl.cell(i, 0), self.truncate_text(cat, 40),
                                   fill=fill, font_color=self.t.TEXT_PRIMARY)
            self._style_table_cell(tbl.cell(i, 1), val, fill=fill, font_color=self.t.TEXT_PRIMARY)

    # ──────────────────────────────────── RIGHT ANALYSIS PANEL ──────────────
    def _render_right_analysis_panel(
        self, slide: Any, res: AskResult,
        left: float, top: float, width: float, height: float,
        chart_pref: str | None,
    ) -> None:
        """Right panel: conclusion card + 3 metric chips + source.
        Mirrors the web ChartSlide right panel."""
        y = top
        pad = 0.12

        # ── Key conclusion card ─────────────────────────────────────────────
        conclusion = ""
        if res.analysis and res.analysis.key_conclusion:
            conclusion = self.truncate_text(res.analysis.key_conclusion, 300)

        card_h = 2.05
        self._card(slide, left, y, width, card_h)
        # Amber lightbulb icon
        self._icon_box(slide, left + pad, y + 0.12, 0.22,
                       self.t.ICONS["bulb"], self.t.AMBER_BG, self.t.AMBER_L)
        self._textbox(slide, left + pad + 0.28, y + 0.10, width - pad * 2 - 0.30, 0.24,
                      "КЛЮЧЕВОЙ ВЫВОД", font_size=7.5, bold=True, color=self.t.TEXT_MUTED)
        self._textbox(slide, left + pad, y + 0.38, width - pad * 2, card_h - 0.50,
                      conclusion, font_size=10, color=self.t.TEXT_SECONDARY, wrap=True)
        y += card_h + 0.10

        # ── 3 metric chips ──────────────────────────────────────────────────
        chips = [
            ("ТРЕНД",         "↑ Рост",    self.t.ICONS["trend_up"], self.t.ACCENTS[2]),  # emerald
            ("АНОМАЛИИ",      "Выявлены",  self.t.ICONS["alert"],    self.t.ACCENTS[3]),  # amber
            ("ДОСТОВЕРНОСТЬ", "Высокая",   self.t.ICONS["check"],    self.t.ACCENTS[1]),  # blue
        ]
        chip_h = 0.60
        for label, value, symbol, ac in chips:
            self._card(slide, left, y, width, chip_h, fill=ac["bg"], border=ac["border"])
            self._icon_box(slide, left + 0.10, y + (chip_h - 0.22) / 2, 0.22,
                           symbol, ac["fg"], self.t.WHITE)
            self._textbox(slide, left + 0.38, y + 0.05, width - 0.48, 0.22,
                          label, font_size=7, color=self.t.TEXT_MUTED)
            self._textbox(slide, left + 0.38, y + 0.24, width - 0.48, 0.32,
                          value, font_size=11, bold=True, color=ac["light"])
            y += chip_h + 0.08

        # ── Source card ─────────────────────────────────────────────────────
        remaining = max(0.42, (top + height) - y)
        self._card(slide, left, y, width, remaining, fill=self.t.CARD_DARK)
        self._textbox(slide, left + pad, y + 0.06, width - pad * 2, 0.20,
                      "Источник данных", font_size=7.5, color=self.t.TEXT_DIMMED)
        self._textbox(slide, left + pad, y + 0.24, width - pad * 2, remaining - 0.30,
                      "Министерство по налогам и сборам РБ",
                      font_size=9.5, bold=True, color=self.t.TEXT_SECONDARY, wrap=True)

    # ──────────────────────────────────────── LAYOUTS ────────────────────────
    def render_slide_header(self, slide: Any, title: str, subtitle: str | None = None) -> None:
        """Draw gradient bar + DATA badge + title + optional subtitle."""
        self._bg(slide)
        self._gradient_bar(slide, 0.0, height=0.07,
                           colors=[self.t.EMERALD, RGBColor(30, 160, 100),
                                   self.t.BLUE, RGBColor(90, 90, 240), self.t.VIOLET])

        title = self.truncate_text(title, 80)
        font_size = 20 if len(title) > 60 else 22
        badge_w = self._badge_tag(slide, self.t.M, 0.27, "DATA",
                                  fill=self.t.EMERALD_BG, border=self.t.EMERALD_B, color=self.t.EMERALD_L)
        self._textbox(slide, self.t.M + badge_w + 0.18, 0.22, self.t.USABLE_W - badge_w - 0.25,
                      0.58, title, font_size=font_size, bold=True)
        if subtitle:
            self._textbox(slide, self.t.M, 0.78, self.t.USABLE_W, 0.28,
                          self.truncate_text(subtitle, 120),
                          font_size=10, color=self.t.TEXT_SECONDARY)

    def _draw_chart_area(
        self, slide: Any, png_path: str | None,
        left: float, top: float, width: float, height: float,
    ) -> None:
        """Draw white chart card + place PNG. If no PNG, draw placeholder."""
        # White card background
        self._rect(slide, left, top, width, height,
                   fill=RGBColor(255, 255, 255), border=self.t.BORDER)
        if png_path:
            ok = self._place_chart_image(slide, png_path, left=left, top=top, width=width, height=height)
        else:
            ok = False
        if not ok:
            # Placeholder text inside the card
            self._textbox(slide, left + 0.5, top + height / 2 - 0.35, width - 1.0, 0.70,
                          "График загружается...",
                          font_size=14, color=RGBColor(150, 150, 150), align=PP_ALIGN.CENTER)

    def _render_top_bottom_layout(
        self, slide: Any, res: AskResult, png_path: str | None, chart_pref: str | None,
    ) -> None:
        m = self.t.M
        chart_h = 3.50
        self._draw_chart_area(slide, png_path, m, self.t.CONTENT_TOP, self.t.USABLE_W, chart_h)
        # Insights row
        if res.analysis and res.analysis.insights:
            insights = res.analysis.insights[:3]
            col_w = (self.t.USABLE_W - (len(insights) - 1) * self.t.GAP) / max(len(insights), 1)
            ins_top = self.t.CONTENT_TOP + chart_h + 0.12
            for ci, insight in enumerate(insights):
                cx = m + ci * (col_w + self.t.GAP)
                ac = self.t.ACCENTS[ci % len(self.t.ACCENTS)]
                self._card(slide, cx, ins_top, col_w, 1.50)
                self._rect(slide, cx, ins_top, col_w, 0.04, fill=ac["fg"], no_line=True)
                self._textbox(slide, cx + 0.12, ins_top + 0.10, col_w - 0.24, 1.36,
                               self.truncate_text(insight, 120),
                               font_size=10, color=self.t.TEXT_SECONDARY, wrap=True)

    def _render_side_by_side_layout(
        self, slide: Any, res: AskResult, png_path: str | None, chart_pref: str | None,
    ) -> None:
        chart_w = 7.20
        right_left = self.t.M + chart_w + self.t.GAP
        right_w = self.t.USABLE_W - chart_w - self.t.GAP
        self._draw_chart_area(slide, png_path, self.t.M, self.t.CONTENT_TOP, chart_w, self.t.CONTENT_H)
        self._render_right_analysis_panel(slide, res, right_left, self.t.CONTENT_TOP,
                                          right_w, self.t.CONTENT_H, chart_pref)

    def _render_chart_and_table_layout(
        self, slide: Any, res: AskResult, png_path: str | None, chart_pref: str | None,
    ) -> None:
        chart_w = 7.20
        right_left = self.t.M + chart_w + self.t.GAP
        right_w = self.t.USABLE_W - chart_w - self.t.GAP
        self._draw_chart_area(slide, png_path, self.t.M, self.t.CONTENT_TOP, chart_w, self.t.CONTENT_H)
        self._render_right_analysis_panel(slide, res, right_left, self.t.CONTENT_TOP,
                                          right_w, self.t.CONTENT_H, chart_pref)

    def _render_kpi_full_layout(
        self, slide: Any, res: AskResult, png_path: str | None, chart_pref: str | None,
    ) -> None:
        chart_w = 8.0
        right_left = self.t.M + chart_w + self.t.GAP
        right_w = self.t.USABLE_W - chart_w - self.t.GAP
        self._draw_chart_area(slide, png_path, self.t.M, self.t.CONTENT_TOP, chart_w, self.t.CONTENT_H)
        self._render_right_analysis_panel(slide, res, right_left, self.t.CONTENT_TOP,
                                          right_w, self.t.CONTENT_H, chart_pref)

    # ──────────────────────────────────────── DECK SLIDES ────────────────────

    def create_title_slide(self, slide: Any, *, date_str: str | None = None) -> None:
        """Title slide — mirrors TitleSlide React component."""
        self._bg(slide)
        # Top gradient bar (violet → blue → cyan → emerald)
        self._gradient_bar(slide, 0.0, height=0.08,
                           colors=[self.t.VIOLET, RGBColor(90, 90, 240), self.t.BLUE,
                                   RGBColor(30, 155, 200), self.t.EMERALD])

        m = self.t.M

        # ── Header row ──────────────────────────────────────────────────────
        # Prototip BI logo icon
        self._icon_box(slide, m, 0.18, 0.40, self.t.ICONS["bar_chart"],
                       self.t.VIOLET_BG, self.t.VIOLET_L)
        self._textbox(slide, m + 0.48, 0.18, 3.2, 0.25, "Prototip BI",
                      font_size=10, bold=True, color=RGBColor(200, 200, 225))
        self._textbox(slide, m + 0.48, 0.40, 3.2, 0.22, "Enterprise AI Analytics",
                      font_size=7.5, color=self.t.TEXT_MUTED)
        # AI-generation badge (top right)
        self._rect(slide, 10.50, 0.22, 1.50, 0.30, fill=self.t.EMERALD_BG, border=self.t.EMERALD_B)
        self._textbox(slide, 10.55, 0.25, 1.38, 0.24, "✓  AI-генерация",
                      font_size=8, color=self.t.EMERALD_L, align=PP_ALIGN.CENTER)
        # Year badge
        self._rect(slide, 12.05, 0.22, 0.85, 0.30, fill=self.t.CARD_BG, border=self.t.BORDER)
        self._textbox(slide, 12.10, 0.25, 0.76, 0.24, str(datetime.now().year),
                      font_size=8, color=self.t.TEXT_MUTED, align=PP_ALIGN.CENTER)

        # ── BI tag ──────────────────────────────────────────────────────────
        self._rect(slide, m, 0.96, 3.80, 0.36, fill=self.t.VIOLET_BG, border=self.t.VIOLET_B)
        self._icon_box(slide, m + 0.12, 1.03, 0.22, "●", self.t.VIOLET, self.t.VIOLET_L)
        self._textbox(slide, m + 0.40, 1.00, 3.36, 0.30,
                      "BI-аналитика  ·  Налоговые поступления РБ",
                      font_size=8.5, bold=True, color=self.t.VIOLET_L)

        # ── Main title ──────────────────────────────────────────────────────
        self._textbox(slide, m, 1.50, 9.6, 1.25,
                      "BI-аналитика налогов РБ",
                      font_size=40, bold=True, color=self.t.WHITE)

        # ── Subtitle ────────────────────────────────────────────────────────
        self._textbox(slide, m, 2.85, 9.6, 0.50,
                      "Синтетические данные (демо), Республика Беларусь",
                      font_size=15, color=self.t.TEXT_SECONDARY)

        # ── Stats grid (4 cards) ─────────────────────────────────────────────
        stats = [
            ("7 регионов",  "Регионов охвачено", self.t.ICONS["pin"],       self.t.ACCENTS[0]),
            ("12+",         "Видов налогов",      self.t.ICONS["file"],      self.t.ACCENTS[1]),
            ("5 агентов",   "AI-система",          self.t.ICONS["cpu"],       self.t.ACCENTS[2]),
            ("20+",         "Инсайтов",            self.t.ICONS["bulb"],      self.t.ACCENTS[3]),
        ]
        card_w = (9.60 - 3 * 0.16) / 4
        cx = m
        for val, lbl, symbol, ac in stats:
            self._card(slide, cx, 3.52, card_w, 1.28)
            self._icon_box(slide, cx + 0.14, 3.66, 0.26, symbol, ac["bg"], ac["light"])
            self._textbox(slide, cx + 0.14, 4.00, card_w - 0.28, 0.42,
                          val, font_size=22, bold=True, color=self.t.WHITE)
            self._textbox(slide, cx + 0.14, 4.44, card_w - 0.28, 0.28,
                          lbl, font_size=8, color=self.t.TEXT_MUTED)
            cx += card_w + 0.16

        # ── Date ─────────────────────────────────────────────────────────────
        date_val = date_str or datetime.now().strftime("%d.%m.%Y")
        self._textbox(slide, m, 4.98, 9.6, 0.40, date_val,
                      font_size=14, color=self.t.WHITE, align=PP_ALIGN.CENTER)

        self._footer_bar(slide)
        self._gov_badge(slide)

    def create_summary_slide(self, slide: Any, overview: str, first_result: AskResult | None) -> None:
        """Summary slide — mirrors SummarySlide React component."""
        self._bg(slide)
        self._gradient_bar(slide, 0.0, height=0.06,
                           colors=[self.t.BLUE, RGBColor(80, 80, 220), self.t.VIOLET,
                                   RGBColor(130, 50, 240), RGBColor(160, 40, 200)])

        m = self.t.M

        # ── Badge + title ────────────────────────────────────────────────────
        badge_w = self._badge_tag(slide, m, 0.22, "SUMMARY",
                                  fill=self.t.BLUE_BG, border=self.t.BLUE_B, color=self.t.BLUE_L)
        self._textbox(slide, m + badge_w + 0.18, 0.18, 10.0, 0.58,
                      "Обзор", font_size=24, bold=True)

        lp_top = 0.90
        left_w = 7.55

        # ── Left panel ───────────────────────────────────────────────────────
        # Main text card
        text_card_h = 3.55
        self._card(slide, m, lp_top, left_w, text_card_h)
        self._icon_box(slide, m + 0.14, lp_top + 0.14, 0.24,
                       self.t.ICONS["file"], self.t.BLUE_BG, self.t.BLUE_L)
        self._textbox(slide, m + 0.44, lp_top + 0.14, 4.0, 0.26,
                      "АНАЛИТИЧЕСКИЙ ОБЗОР", font_size=7.5, bold=True, color=self.t.TEXT_MUTED)
        self._textbox(slide, m + 0.14, lp_top + 0.46, left_w - 0.28, text_card_h - 0.58,
                      self.truncate_text(overview, 580),
                      font_size=11, color=self.t.TEXT_SECONDARY, wrap=True)

        # 2 info chips below
        info_top = lp_top + text_card_h + 0.12
        info_w = (left_w - self.t.GAP) / 2
        info_data = [
            ("Отчётный период", "2024 год",      self.t.ICONS["file"],   self.t.ACCENTS[1]),
            ("Источник данных",  "ClickHouse OLAP", self.t.ICONS["activity"], self.t.ACCENTS[2]),
        ]
        for ci, (lbl, val, symbol, ac) in enumerate(info_data):
            cx = m + ci * (info_w + self.t.GAP)
            self._card(slide, cx, info_top, info_w, 0.78)
            self._icon_box(slide, cx + 0.12, info_top + 0.22, 0.28,
                           symbol, self.t.CARD_BG2, ac["light"])
            self._textbox(slide, cx + 0.48, info_top + 0.08, info_w - 0.58, 0.22,
                          lbl, font_size=7.5, color=self.t.TEXT_MUTED)
            self._textbox(slide, cx + 0.48, info_top + 0.30, info_w - 0.58, 0.38,
                          val, font_size=11, bold=True)

        # ── Right panel ──────────────────────────────────────────────────────
        right_left = m + left_w + self.t.GAP
        right_w = self.t.USABLE_W - left_w - self.t.GAP
        rp_top = lp_top

        self._textbox(slide, right_left, rp_top, right_w, 0.22, "КЛЮЧЕВЫЕ ПАРАМЕТРЫ",
                      font_size=7.5, bold=True, color=self.t.TEXT_MUTED)

        kpi_params = [
            ("7 регионов",  "Охват",       self.t.ICONS["pin"],      self.t.ACCENTS[0]),
            ("Text-to-SQL", "Методология", self.t.ICONS["cpu"],      self.t.ACCENTS[1]),
            ("98.5%",       "Точность",    self.t.ICONS["target"],   self.t.ACCENTS[2]),
            ("5 агентов",   "AI-система",  self.t.ICONS["activity"], self.t.ACCENTS[3]),
        ]
        chip_h = 0.70
        ry = rp_top + 0.24
        for val, lbl, symbol, ac in kpi_params:
            self._card(slide, right_left, ry, right_w, chip_h, fill=ac["bg"], border=ac["border"])
            self._icon_box(slide, right_left + 0.12, ry + (chip_h - 0.28) / 2, 0.28,
                           symbol, self.t.CARD_BG, ac["light"])
            self._textbox(slide, right_left + 0.48, ry + 0.08, right_w - 0.58, 0.22,
                          lbl, font_size=7.5, color=self.t.TEXT_MUTED)
            self._textbox(slide, right_left + 0.48, ry + 0.30, right_w - 0.58, chip_h - 0.40,
                          val, font_size=11, bold=True, color=ac["light"])
            # Chevron
            self._textbox(slide, right_left + right_w - 0.28, ry + 0.24, 0.22, 0.26,
                          "›", font_size=14, color=self.t.TEXT_DIMMED, align=PP_ALIGN.RIGHT)
            ry += chip_h + 0.10

        # Trend card
        trend_top = ry
        trend_h = max(0.55, (lp_top + text_card_h + 0.12 + 0.78) - trend_top)
        self._card(slide, right_left, trend_top, right_w, trend_h,
                   fill=self.t.EMERALD_BG, border=self.t.EMERALD_B)
        self._icon_box(slide, right_left + 0.12, trend_top + 0.10, 0.22,
                       self.t.ICONS["trend_up"], self.t.EMERALD, self.t.WHITE)
        self._textbox(slide, right_left + 0.40, trend_top + 0.08, right_w - 0.50, 0.24,
                      "Тренд", font_size=9, bold=True, color=self.t.EMERALD_L)
        self._textbox(slide, right_left + 0.40, trend_top + 0.30, right_w - 0.50,
                      max(0.22, trend_h - 0.36),
                      "Рост налоговых поступлений в анализируемом периоде",
                      font_size=8.5, color=self.t.TEXT_SECONDARY, wrap=True)

        self._footer_bar(slide)
        self._gov_badge(slide)

    def create_themes_slide(self, slide: Any, themes: list[str], questions: list[str]) -> None:
        """Themes / Agenda slide — mirrors ThemesSlide React component."""
        self._bg(slide)
        self._gradient_bar(slide, 0.0, height=0.06,
                           colors=[self.t.VIOLET, RGBColor(80, 80, 240), self.t.BLUE,
                                   RGBColor(20, 160, 210), self.t.CYAN])

        m = self.t.M

        # Badge + title
        badge_w = self._badge_tag(slide, m, 0.22, "AGENDA",
                                  fill=self.t.VIOLET_BG, border=self.t.VIOLET_B, color=self.t.VIOLET_L)
        self._textbox(slide, m + badge_w + 0.18, 0.18, 9.5, 0.58,
                      "Темы и повестка", font_size=24, bold=True)
        self._textbox(slide, 11.5, 0.26, 1.5, 0.32, f"{len(themes)} разделов",
                      font_size=9, color=self.t.TEXT_MUTED, align=PP_ALIGN.RIGHT)

        # Theme cards
        themes_top = 0.95
        n = min(7, len(themes))
        avail_h = 5.75
        card_h = min(0.78, (avail_h - (n - 1) * 0.10) / max(n, 1))

        # Icon symbols cycling
        theme_icons = [
            self.t.ICONS["target"], self.t.ICONS["globe"],  self.t.ICONS["bar_chart"],
            self.t.ICONS["zap"],    self.t.ICONS["shield"],  self.t.ICONS["cpu"],
            self.t.ICONS["activity"],
        ]

        cy = themes_top
        for i, theme in enumerate(themes[:n]):
            ac = self.t.ACCENTS[i % len(self.t.ACCENTS)]
            self._card(slide, m, cy, self.t.USABLE_W, card_h, fill=ac["bg"], border=ac["border"])
            # Accent top strip
            self._rect(slide, m, cy, self.t.USABLE_W, 0.04, fill=ac["fg"], no_line=True)
            # Icon
            self._icon_box(slide, m + 0.12, cy + (card_h - 0.32) / 2, 0.32,
                           theme_icons[i % len(theme_icons)], self.t.CARD_BG, ac["light"])
            # Section label + text
            self._textbox(slide, m + 0.54, cy + 0.04, 2.0, 0.20,
                          f"Раздел {i + 1}", font_size=7, color=self.t.TEXT_MUTED)
            self._textbox(slide, m + 0.54, cy + 0.22, self.t.USABLE_W - 1.30, card_h - 0.30,
                          self.truncate_text(theme, 100), font_size=11, color=self.t.TEXT_PRIMARY)
            # Numbered badge
            nbw = 0.34
            self._rect(slide, m + self.t.USABLE_W - nbw - 0.12, cy + (card_h - 0.28) / 2,
                       nbw, 0.28, fill=ac["fg"], no_line=True)
            num_str = f"0{i + 1}" if i < 9 else str(i + 1)
            self._textbox(slide, m + self.t.USABLE_W - nbw - 0.12, cy + (card_h - 0.28) / 2 + 0.02,
                          nbw, 0.26, num_str, font_size=9, bold=True, color=self.t.WHITE,
                          align=PP_ALIGN.CENTER)
            cy += card_h + 0.10

        # Bottom info row
        info_top = max(cy + 0.04, 6.28)
        info_items = [
            ("Методология", "AI Text-to-SQL + BI", self.t.ICONS["cpu"]),
            ("Источник",    "ClickHouse OLAP",      self.t.ICONS["activity"]),
            ("Охват",       "Все регионы РБ",       self.t.ICONS["globe"]),
        ]
        chip_w = (self.t.USABLE_W - 2 * self.t.GAP) / 3
        for ci, (lbl, val, symbol) in enumerate(info_items):
            cx = m + ci * (chip_w + self.t.GAP)
            self._card(slide, cx, info_top, chip_w, 0.52)
            ac = self.t.ACCENTS[ci + 1]
            self._icon_box(slide, cx + 0.12, info_top + 0.14, 0.22, symbol, self.t.CARD_BG2, ac["light"])
            self._textbox(slide, cx + 0.40, info_top + 0.06, chip_w - 0.50, 0.20,
                          lbl, font_size=7.5, color=self.t.TEXT_MUTED)
            self._textbox(slide, cx + 0.40, info_top + 0.26, chip_w - 0.50, 0.22,
                          val, font_size=9, bold=True, color=self.t.TEXT_SECONDARY)

        self._footer_bar(slide)
        self._gov_badge(slide)

    def render_question_slide(
        self, slide: Any, question: str, res: AskResult,
        chart_pref: str | None, png_path: str | None,
    ) -> None:
        """Chart/data slide — mirrors ChartSlide React component."""
        header, subtitle = self._resolve_headers(question, res)
        self.render_slide_header(slide, header, subtitle)

        ctype = chart_pref or (getattr(res.chart_spec, "chart_type", None) if res.chart_spec else None)
        layout = self.choose_layout(ctype)

        if layout == LayoutKind.TOP_BOTTOM:
            self._render_top_bottom_layout(slide, res, png_path, chart_pref)
        elif layout == LayoutKind.SIDE_BY_SIDE:
            self._render_side_by_side_layout(slide, res, png_path, chart_pref)
        elif layout == LayoutKind.CHART_TABLE:
            self._render_chart_and_table_layout(slide, res, png_path, chart_pref)
        else:
            self._render_kpi_full_layout(slide, res, png_path, chart_pref)
        # NOTE: footer/badge is drawn by build_question_slides_footer() called externally

    def create_takeaways_slide(self, slide: Any, takeaways: list[str]) -> None:
        """Takeaways slide — mirrors TakeawaysSlide React component."""
        self._bg(slide)
        self._gradient_bar(slide, 0.0, height=0.06,
                           colors=[self.t.EMERALD, RGBColor(20, 170, 110),
                                   RGBColor(20, 165, 150), RGBColor(10, 170, 190), self.t.CYAN])

        m = self.t.M
        shown = takeaways[:10]
        n = len(shown)

        badge_w = self._badge_tag(slide, m, 0.22, "KEY INSIGHTS",
                                  fill=self.t.EMERALD_BG, border=self.t.EMERALD_B, color=self.t.EMERALD_L)
        self._textbox(slide, m + badge_w + 0.18, 0.18, 9.5, 0.58,
                      "Ключевые выводы", font_size=24, bold=True)
        self._textbox(slide, 11.0, 0.28, 2.1, 0.30,
                      f"{n} ключевых вывода", font_size=9, color=self.t.TEXT_MUTED,
                      align=PP_ALIGN.RIGHT)

        grid_top = 0.90
        use_grid = n > 5
        cols = 2 if use_grid else 1
        col_w = (self.t.USABLE_W - (cols - 1) * self.t.GAP) / cols
        rows = (n + cols - 1) // cols
        avail_h = 5.70
        card_h = min(1.12, (avail_h - (rows - 1) * 0.10) / max(rows, 1))

        for i, tw in enumerate(shown):
            col = i % cols
            row = i // cols
            cx = m + col * (col_w + self.t.GAP)
            cy = grid_top + row * (card_h + 0.10)
            ac = self.t.ACCENTS[i % len(self.t.ACCENTS)]

            self._card(slide, cx, cy, col_w, card_h)
            # Colored number badge
            num_sz = 0.30
            self._icon_box(slide, cx + 0.12, cy + 0.12, num_sz,
                           str(i + 1), ac["fg"], self.t.WHITE)
            # Text content
            text_h = card_h - 0.24 - 0.12
            max_chars = 110 if use_grid else 190
            self._textbox(slide, cx + 0.52, cy + 0.10, col_w - 0.64, text_h,
                          self.truncate_text(tw, max_chars),
                          font_size=10 if use_grid else 12,
                          color=self.t.TEXT_SECONDARY, wrap=True)
            # Progress bar
            pct = max(0.40, 1.0 - i * 0.055)
            self._progress_bar(slide, cx + 0.52, cy + card_h - 0.12, col_w - 0.64, pct, ac)

        # Bottom confirmation bar
        bottom_top = grid_top + rows * (card_h + 0.10) + 0.04
        if bottom_top > 6.52:
            bottom_top = 6.52
        self._card(slide, m, bottom_top, self.t.USABLE_W, 0.52,
                   fill=self.t.EMERALD_BG, border=self.t.EMERALD_B)
        self._icon_box(slide, m + 0.14, bottom_top + 0.14, 0.24,
                       self.t.ICONS["check"], self.t.EMERALD, self.t.WHITE)
        self._textbox(slide, m + 0.46, bottom_top + 0.09, self.t.USABLE_W - 0.60, 0.38,
                      f"{n} выводов сформированы на основе AI-анализа налоговых данных "
                      "Республики Беларусь (синтетический датасет)",
                      font_size=9, color=self.t.TEXT_SECONDARY, wrap=True)

        self._footer_bar(slide)
        self._gov_badge(slide)

    def create_recommendations_slide(self, slide: Any, recommendations: list[str]) -> None:
        """Recommendations slide — mirrors RecommendationsSlide React component."""
        self._bg(slide)
        self._gradient_bar(slide, 0.0, height=0.06,
                           colors=[self.t.AMBER, RGBColor(230, 120, 20), RGBColor(240, 80, 40),
                                   RGBColor(240, 50, 70), self.t.ROSE])

        m = self.t.M
        recs = recommendations[:4]

        badge_w = self._badge_tag(slide, m, 0.22, "ACTIONS",
                                  fill=self.t.AMBER_BG, border=self.t.AMBER_B, color=self.t.AMBER_L)
        self._textbox(slide, m + badge_w + 0.18, 0.18, 9.5, 0.58,
                      "Рекомендации", font_size=24, bold=True)
        # Count badge
        cnt_txt = f"{len(recs)} рекомендаций"
        cnt_w = max(1.3, len(cnt_txt) * 0.10 + 0.2)
        self._rect(slide, 12.5 - cnt_w, 0.22, cnt_w, 0.30, fill=self.t.AMBER_BG, border=self.t.AMBER_B)
        self._textbox(slide, 12.5 - cnt_w + 0.06, 0.24, cnt_w - 0.12, 0.26,
                      cnt_txt, font_size=8, bold=True, color=self.t.AMBER_L, align=PP_ALIGN.CENTER)

        # Rec icons cycling
        rec_icons = [
            self.t.ICONS["arrow"], self.t.ICONS["zap"], self.t.ICONS["target"],
            self.t.ICONS["shield"],
        ]

        grid_top = 0.90
        grid_h = 5.30
        col_w = (self.t.USABLE_W - self.t.GAP) / 2
        row_h = (grid_h - self.t.GAP) / 2

        card_accents = [self.t.ACCENTS[0], self.t.ACCENTS[1], self.t.ACCENTS[2], self.t.ACCENTS[3]]
        for i, rec in enumerate(recs):
            col = i % 2
            row = i // 2
            cx = m + col * (col_w + self.t.GAP)
            cy = grid_top + row * (row_h + self.t.GAP)
            ac = card_accents[i]

            self._card(slide, cx, cy, col_w, row_h, fill=ac["bg"], border=ac["border"])
            # Top accent strip
            self._rect(slide, cx, cy, col_w, 0.05, fill=ac["fg"], no_line=True)

            # Icon button
            self._icon_box(slide, cx + 0.14, cy + 0.14, 0.36,
                           rec_icons[i % len(rec_icons)], ac["fg"], self.t.WHITE)

            # Priority label
            self._textbox(slide, cx + col_w - 1.15, cy + 0.18, 1.05, 0.24,
                          f"ПРИОРИТЕТ {i + 1}",
                          font_size=7.5, bold=True, color=self.t.TEXT_MUTED, align=PP_ALIGN.RIGHT)

            # Rec text
            self._textbox(slide, cx + 0.14, cy + 0.60, col_w - 0.28, row_h - 1.00,
                          self.truncate_text(rec, 160),
                          font_size=10.5, color=self.t.TEXT_SECONDARY, wrap=True)

            # Progress bar + pct
            pct = max(0.30, 0.90 - i * 0.15)
            bar_top = cy + row_h - 0.30
            self._progress_bar(slide, cx + 0.14, bar_top, col_w - 0.62, pct, ac)
            self._textbox(slide, cx + col_w - 0.52, bar_top - 0.06, 0.44, 0.22,
                          f"{int(pct * 100)}%",
                          font_size=8.5, bold=True, color=ac["light"], align=PP_ALIGN.RIGHT)

        # Disclaimer bar
        disc_top = grid_top + grid_h + 0.10
        if disc_top > 6.52:
            disc_top = 6.52
        self._card(slide, m, disc_top, self.t.USABLE_W, 0.52)
        self._icon_box(slide, m + 0.14, disc_top + 0.14, 0.24,
                       self.t.ICONS["shield"], self.t.AMBER, self.t.WHITE)
        self._textbox(slide, m + 0.46, disc_top + 0.09, self.t.USABLE_W - 0.60, 0.40,
                      "Рекомендации сформированы AI-системой. "
                      "Требуют экспертной валидации перед внедрением.",
                      font_size=9, color=self.t.TEXT_SECONDARY, wrap=True)

        self._footer_bar(slide)
        self._gov_badge(slide)

    def create_appendix_slide(self, slide: Any, slide_num: int) -> None:
        self._bg(slide)
        self._gradient_bar(slide, 0.0)
        self._textbox(slide, 0.5, 2.8, 12.333, 1.0, "Приложение",
                      font_size=40, bold=True, color=self.t.TEXT_PRIMARY, align=PP_ALIGN.CENTER)
        self._textbox(slide, 0.5, 4.0, 12.333, 0.60,
                      "Дополнительные материалы и методология расчетов",
                      font_size=20, color=self.t.TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        self._footer_bar(slide, slide_num=slide_num)

    def build_question_slides_footer(self, slide: Any, slide_num: int) -> None:
        """Called after render_question_slide to add footer with slide number."""
        self._footer_bar(slide, slide_num=slide_num)
        self._gov_badge(slide)

    # ── Backward-compat shims ─────────────────────────────────────────────────
    def _apply_dark_bg(self, slide: Any) -> None:
        self._bg(slide)

    def _add_footer(self, slide: Any, *, slide_num: int | None = None) -> None:
        self._footer_bar(slide, slide_num=slide_num)

    def _add_gov_badge(self, slide: Any) -> None:
        self._gov_badge(slide)

    def _add_centered_text(
        self, slide: Any, text: str, *, top: Any,
        font_size: int = 32, bold: bool = False, color: RGBColor | None = None,
    ) -> None:
        try:
            top_in = top.inches
        except AttributeError:
            top_in = float(top) / 914400.0
        self._textbox(slide, 0.5, top_in, 12.333, 1.2, text,
                      font_size=font_size, bold=bold,
                      color=color or self.t.TEXT_PRIMARY, align=PP_ALIGN.CENTER)

    def _add_title_text(
        self, slide: Any, text: str, *, top: float,
        font_size: int = 26, bold: bool = True,
        color: RGBColor | None = None,
        left: float | None = None, width: float | None = None,
    ) -> None:
        self._textbox(slide, left or self.t.M, top, width or self.t.USABLE_W, 0.8,
                      text, font_size=font_size, bold=bold, color=color or self.t.TEXT_PRIMARY)

    def _add_insights_block(
        self, slide: Any, res: AskResult, *,
        left: float, top: float, width: float, height: float,
        chart_pref: str | None = None, columns: int = 1,
    ) -> None:
        self._render_right_analysis_panel(slide, res, left, top, width, height, chart_pref)

    def _draw_kpi_card(
        self, slide: Any, left: float, top: float, width: float, height: float,
        title: str, value: str,
    ) -> None:
        self._card(slide, left, top, width, height)
        self._textbox(slide, left + 0.12, top + 0.10, width - 0.24, 0.28,
                      title, font_size=10, color=self.t.TEXT_SECONDARY)
        self._textbox(slide, left + 0.12, top + 0.38, width - 0.24, height - 0.50,
                      value, font_size=20, bold=True, color=self.t.BLUE)